"""
AI瓒呰〃闈㈢粨鏋勮壊鏅鸿兘璁捐绯荤粺 - 瀛︽湳绛旇京PPT鐢熸垚鍣?
闀挎矙鐞嗗伐澶у 鐗╃悊涓庣數瀛愮瀛﹀闄?鍏夌數2501鐝?
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# ==================== 鍏ㄥ眬閰嶇疆 ====================

# 閰嶈壊浣撶郴
COLOR_PRIMARY = RGBColor(0x16, 0x5D, 0xFF)    # 娣辩┖钃?#165DFF
COLOR_TEXT = RGBColor(0x1D, 0x21, 0x29)        # 鐐伆鑹?#1D2129
COLOR_BG = RGBColor(0xF7, 0xF8, 0xFA)          # 娴呯伆鐧?#F7F8FA
COLOR_ACCENT = RGBColor(0xFF, 0x7D, 0x00)      # 鐞ョ弨閲?#FF7D00

# 瀛椾綋閰嶇疆锛堣纭繚绯荤粺瀹夎鎬濇簮榛戜綋锛屾垨鏇挎崲涓篗icrosoft YaHei锛?
FONT_TITLE = "鎬濇簮榛戜綋 Bold"
FONT_BODY = "鎬濇簮榛戜綋 Regular"
FONT_MATH = "Times New Roman"

# 椤甸潰灏哄锛?6:9锛?
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def create_presentation():
    """鍒涘缓婕旂ず鏂囩"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # 璁剧疆榛樿鑳屾櫙鑹?
    for layout in prs.slide_layouts:
        layout.background.fill.solid()
        layout.background.fill.fore_color.rgb = COLOR_BG
    
    return prs

def add_title_shape(slide, text, left, top, width, height, 
                    font_size=Pt(44), bold=True, color=COLOR_PRIMARY,
                    font_name=FONT_TITLE):
    """娣诲姞鏍囬鏂囨湰妗?""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = PP_ALIGN.LEFT
    return shape

def add_body_text(slide, text, left, top, width, height,
                  font_size=Pt(14), color=COLOR_TEXT,
                  font_name=FONT_BODY, align=PP_ALIGN.LEFT):
    """娣诲姞姝ｆ枃鏂囨湰妗?""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return shape

def add_formula_text(slide, text, left, top, width, height,
                     font_size=Pt(12), color=COLOR_TEXT):
    """娣诲姞鍏紡鏂囨湰锛圱imes New Roman锛屾枩浣撳彉閲忥級"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    
    # 瑙ｆ瀽鍏紡鏂囨湰锛屽鐞嗘枩浣?
    # 绠€鍖栧鐞嗭細鏁翠綋浣跨敤Times New Roman锛屽叧閿彉閲忔墜鍔ㄦ爣璁版枩浣?
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = FONT_MATH
    p.alignment = PP_ALIGN.LEFT
    
    return shape

def add_rectangle(slide, left, top, width, height, 
                  fill_color=None, line_color=None, line_width=Pt(1)):
    """娣诲姞鐭╁舰褰㈢姸"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    
    return shape

def add_line(slide, x1, y1, x2, y2, color=COLOR_PRIMARY, width=Pt(2)):
    """娣诲姞绾挎潯"""
    shape = slide.shapes.add_connector(
        MSO_SHAPE.STRAIGHT_CONNECTOR_1, x1, y1, x2, y2
    )
    shape.line.color.rgb = color
    shape.line.width = width
    return shape

# ==================== 绗?椤碉細灏侀潰椤?====================

def create_slide_1_cover(prs):
    """灏侀潰椤?""
    slide_layout = prs.slide_layouts[6]  # 绌虹櫧甯冨眬
    slide = prs.slides.add_slide(slide_layout)
    
    # 璁剧疆鑳屾櫙
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 宸︿晶60%鏂囧瓧鍖?
    left_margin = Inches(0.8)
    top_margin = Inches(2.0)
    
    # 涓绘爣棰?
    add_title_shape(slide, "AI 瓒呰〃闈㈢粨鏋勮壊鏅鸿兘璁捐绯荤粺",
                    left_margin, Inches(1.8), Inches(7.0), Inches(1.2),
                    font_size=Pt(48), bold=True, color=COLOR_PRIMARY)
    
    # 鍓爣棰?
    add_body_text(slide, "鍩轰簬鐗╃悊妯″瀷涓庢繁搴﹀涔犵殑绾崇背鍏夊蹇€熻璁″钩鍙?,
                  left_margin, Inches(3.2), Inches(7.0), Inches(0.6),
                  font_size=Pt(18), color=COLOR_TEXT)
    
    # 鍒嗛殧绾?
    add_line(slide, left_margin, Inches(4.2), Inches(6.5), Inches(4.2),
             color=COLOR_PRIMARY, width=Pt(2))
    
    # 搴曢儴淇℃伅
    info_text = "闀挎矙鐞嗗伐澶у 路 鐗╃悊涓庣數瀛愮瀛﹀闄?路 鍏夌數2501鐝璡n姹囨姤浜猴細涔斿畨鐞猏n鍥㈤槦鎴愬憳锛氳阿瀹剁彏銆佷警鐞?
    add_body_text(slide, info_text,
                  left_margin, Inches(4.6), Inches(7.0), Inches(1.8),
                  font_size=Pt(14), color=COLOR_TEXT)
    
    # 鍙充晶40%瑙嗚鍖?- 绾崇背鏌遍樀鍒楃ず鎰忓浘鍗犱綅
    # 娣诲姞瑁呴グ鎬х煩褰㈡琛ㄧず鍥惧儚鍖哄煙
    img_placeholder = add_rectangle(slide, Inches(8.5), Inches(2.0), 
                                   Inches(4.0), Inches(4.0),
                                   fill_color=RGBColor(0xE8, 0xEA, 0xED),
                                   line_color=COLOR_PRIMARY, line_width=Pt(1))
    
    # 绀烘剰鍥捐鏄庢枃瀛?
    add_body_text(slide, "[绾崇背鏌遍樀鍒椾刊瑙嗙ず鎰忓浘]\n浣庨ケ鍜屽害绉戝鍙鍖?,
                  Inches(8.5), Inches(6.2), Inches(4.0), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 鍙充笅瑙掍簯绔湴鍧€
    add_body_text(slide, "绯荤粺浜戠鍦板潃锛歨ttps://metasurface-web-2you8jsy6wxhmhqgc9vwhm.streamlit.app\n[浜岀淮鐮佸崰浣峕",
                  Inches(9.5), Inches(6.8), Inches(3.5), Inches(0.6),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.RIGHT)
    
    return slide

# ==================== 绗?椤碉細鐩綍椤?====================

def create_slide_2_toc(prs):
    """鐩綍椤?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰鏍囬
    add_title_shape(slide, "鐩綍",
                    Inches(0.8), Inches(0.6), Inches(2.0), Inches(0.8),
                    font_size=Pt(36), bold=True, color=COLOR_PRIMARY)
    
    # 鍏釜妯″潡 - 绔栧悜鎺掑垪
    modules = [
        ("01", "鐮旂┒鑳屾櫙涓庢剰涔?),
        ("02", "绯荤粺鎬讳綋鏋舵瀯"),
        ("03", "鏍稿績鎶€鏈紩鎿?),
        ("04", "鏍稿績鍔熻兘浣撶郴"),
        ("05", "鏉愭枡浣撶郴涓庢€ц兘楠岃瘉"),
        ("06", "鎬荤粨涓庢湭鏉ュ睍鏈?)
    ]
    
    start_y = Inches(1.6)
    for i, (num, title) in enumerate(modules):
        y_pos = start_y + i * Inches(0.95)
        
        # 搴忓彿锛堟斁澶э紝涓昏壊锛?
        add_title_shape(slide, num,
                       Inches(1.0), y_pos, Inches(0.8), Inches(0.7),
                       font_size=Pt(32), bold=True, color=COLOR_PRIMARY)
        
        # 鏍囬鏂囧瓧
        add_body_text(slide, title,
                     Inches(2.0), y_pos + Inches(0.1), Inches(4.0), Inches(0.6),
                     font_size=Pt(20), color=COLOR_TEXT)
        
        # 鏋佺畝鍥炬爣鍗犱綅锛堝渾褰級
        add_rectangle(slide, Inches(6.5), y_pos + Inches(0.15), 
                     Inches(0.4), Inches(0.4),
                     fill_color=COLOR_PRIMARY, line_color=None)
        
        # 杩炴帴绾?
        if i < len(modules) - 1:
            add_line(slide, Inches(1.4), y_pos + Inches(0.75), 
                    Inches(1.4), y_pos + Inches(0.95),
                    color=RGBColor(0xC9, 0xCD, 0xD4), width=Pt(1))
    
    return slide

# ==================== 绗?椤碉細鐮旂┒鑳屾櫙 ====================

def create_slide_3_background(prs):
    """鐮旂┒鑳屾櫙涓庝紶缁熻璁＄摱棰?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "鐮旂┒鑳屾櫙涓庤涓氱棝鐐?,
                    Inches(0.8), Inches(0.5), Inches(10.0), Inches(0.8),
                    font_size=Pt(32), bold=True, color=COLOR_PRIMARY)
    
    # 涓婂崐娈靛紩瑷€
    intro_text = ("瓒呰〃闈㈢粨鏋勮壊鏄撼绫冲厜瀛愬鐨勬牳蹇冪爺绌舵柟鍚戯紝渚濋潬浜氭尝闀跨粨鏋勪笌鍏夌殑鐩镐簰浣滅敤浜х敓棰滆壊锛?
                 "鍏锋湁楂樺垎杈ㄧ巼銆佹姉瑜壊銆佺幆淇濈瓑浼樺娍锛屽湪楂樻竻鏄剧ず銆佸厜瀛﹂槻浼€佺敓鐗╀紶鎰熴€佸厜浼忎紭鍖栫瓑"
                 "棰嗗煙鍏峰閲嶅ぇ搴旂敤浠峰€笺€俓n\n"
                 "浣嗕紶缁熻秴琛ㄩ潰璁捐渚濊禆鍏ㄦ尝鏁板€间豢鐪燂紝瀛樺湪涓夊ぇ琛屼笟鍏辨€х摱棰堬細")
    add_body_text(slide, intro_text,
                  Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.6),
                  font_size=Pt(14), color=COLOR_TEXT)
    
    # 涓夊垪鐥涚偣鍗＄墖
    card_width = Inches(3.6)
    card_height = Inches(2.4)
    card_y = Inches(3.2)
    cards = [
        ("璁捐鍛ㄦ湡鏋侀暱", 
         "鍗曠粨鏋?FDTD/RCWA 浠跨湡闇€鏁板皬鏃讹紝鍙傛暟绌洪棿閬嶅巻鎴愭湰鏋侀珮锛屾棤娉曞疄鐜板揩閫熻凯浠?),
        ("浣跨敤闂ㄦ楂樻槀", 
         "渚濊禆 Lumerical銆丆OMSOL 绛夊晢涓氫笓涓氳蒋浠讹紝鎺堟潈鎴愭湰楂橈紝瀛︿範鍛ㄦ湡闀胯揪鏁版湀"),
        ("閫嗚璁℃晥鐜囦綆涓?, 
         "鎵嬪姩璋冨弬渚濊禆璁捐鑰呯粡楠岋紝闅句互浠庣洰鏍囬鑹插弽鍚戞帹瀵兼渶浼樼粨鏋勫弬鏁?)
    ]
    
    for i, (title, content) in enumerate(cards):
        x_pos = Inches(0.8) + i * (card_width + Inches(0.4))
        
        # 鍗＄墖鑳屾櫙
        add_rectangle(slide, x_pos, card_y, card_width, card_height,
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(1.5))
        
        # 鍥炬爣鍗犱綅锛堥《閮級
        icon_y = card_y + Inches(0.2)
        add_rectangle(slide, x_pos + Inches(1.4), icon_y, 
                     Inches(0.8), Inches(0.5),
                     fill_color=COLOR_PRIMARY, line_color=None)
        
        # 鏍囬
        add_title_shape(slide, title,
                       x_pos + Inches(0.3), icon_y + Inches(0.6), 
                       card_width - Inches(0.6), Inches(0.5),
                       font_size=Pt(16), bold=True, color=COLOR_PRIMARY)
        
        # 鍐呭
        add_body_text(slide, content,
                     x_pos + Inches(0.3), icon_y + Inches(1.1), 
                     card_width - Inches(0.6), Inches(1.0),
                     font_size=Pt(12), color=COLOR_TEXT)
    
    # 搴曢儴鎬荤粨鍙ワ紙楂樹寒锛?
    summary = "鏈郴缁熸棬鍦ㄦ墦閫?\"闆堕棬妲涖€佹绉掔骇銆佹櫤鑳藉寲\" 鐨勮秴琛ㄩ潰缁撴瀯鑹茶璁″钩鍙?
    add_title_shape(slide, summary,
                    Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6),
                    font_size=Pt(16), bold=True, color=COLOR_ACCENT)
    
    return slide

# ==================== 绗?椤碉細绯荤粺鎬讳綋鏋舵瀯 ====================

def create_slide_4_architecture(prs):
    """绯荤粺鎬讳綋鏋舵瀯"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "绯荤粺鎬讳綋鏋舵瀯涓庢牳蹇冮€昏緫",
                    Inches(0.8), Inches(0.5), Inches(10.0), Inches(0.8),
                    font_size=Pt(32), bold=True, color=COLOR_PRIMARY)
    
    # 涓棿妯悜娴佺▼鍥?
    flow_y = Inches(1.8)
    modules = [
        ("鍙傛暟杈撳叆灞?, "绾崇背鏌辩洿寰?D銆侀珮搴?H銆佸懆鏈?P\n鏉愭枡閫夋嫨銆佽娴嬭搴︺€佸亸鎸ā寮?),
        ("鍙屽紩鎿庤绠楀眰", "鐗╃悊寮曟搸锛堟礇浼﹀吂瑙ｆ瀽妯″瀷锛塡n+ ML 寮曟搸锛圧esMLP 浠ｇ悊妯″瀷锛?),
        ("缁撴灉杈撳嚭灞?, "鍙嶅皠鍏夎氨鏇茬嚎銆丆IE 鑹插潗鏍嘰nRGB 瀹炶壊銆佽繙鍦鸿璋卞垎甯?),
        ("鏅鸿兘鍒嗘瀽灞?, "澶氱畻娉曢€嗚璁°€佸ぇ妯″瀷鐗╃悊瑙ｆ瀽\n宸ヨ壓瀹瑰樊璇勪及")
    ]
    
    box_width = Inches(2.6)
    box_height = Inches(1.4)
    start_x = Inches(0.9)
    
    for i, (title, content) in enumerate(modules):
        x_pos = start_x + i * (box_width + Inches(0.5))
        
        # 妯″潡妗?
        add_rectangle(slide, x_pos, flow_y, box_width, box_height,
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(2))
        
        # 鏍囬
        add_title_shape(slide, title,
                       x_pos + Inches(0.15), flow_y + Inches(0.15), 
                       box_width - Inches(0.3), Inches(0.4),
                       font_size=Pt(14), bold=True, color=COLOR_PRIMARY)
        
        # 鍐呭
        add_body_text(slide, content,
                     x_pos + Inches(0.15), flow_y + Inches(0.55), 
                     box_width - Inches(0.3), Inches(0.8),
                     font_size=Pt(11), color=COLOR_TEXT)
        
        # 绠ご杩炴帴
        if i < len(modules) - 1:
            arrow = add_line(slide, x_pos + box_width, flow_y + box_height/2,
                           x_pos + box_width + Inches(0.5), flow_y + box_height/2,
                           color=COLOR_PRIMARY, width=Pt(3))
    
    # 搴曢儴鍥涘垪鍒涙柊鐐瑰崱鐗?
    card_y = Inches(3.8)
    card_width = Inches(2.9)
    card_height = Inches(2.6)
    
    innovations = [
        "鐗╃悊-娣卞害瀛︿範鍙岄┍鍔紝鍏奸【璁＄畻绮惧害涓庨€熷害",
        "鍥涚閫嗚璁＄畻娉曪紝瑕嗙洊鍏ㄥ満鏅璁￠渶姹?,
        "澶фā鍨嬪師鐢熸帴鍏ワ紝鑷姩鐢熸垚鐗╃悊瑙ｉ噴涓庝紭鍖栧缓璁?,
        "绾?Web 浜戠閮ㄧ讲锛岄浂瀹夎寮€绠卞嵆鐢?
    ]
    
    add_title_shape(slide, "鍥涘ぇ鏍稿績鍒涙柊",
                    Inches(0.8), Inches(3.5), Inches(4.0), Inches(0.4),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    for i, text in enumerate(innovations):
        x_pos = Inches(0.8) + i * (card_width + Inches(0.25))
        
        # 鍗＄墖鑳屾櫙
        add_rectangle(slide, x_pos, card_y, card_width, card_height,
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_ACCENT, line_width=Pt(2))
        
        # 搴忓彿
        add_title_shape(slide, f"{i+1}",
                       x_pos + Inches(0.2), card_y + Inches(0.2), 
                       Inches(0.5), Inches(0.5),
                       font_size=Pt(24), bold=True, color=COLOR_ACCENT)
        
        # 鍐呭
        add_body_text(slide, text,
                     x_pos + Inches(0.2), card_y + Inches(0.8), 
                     card_width - Inches(0.4), Inches(1.6),
                     font_size=Pt(13), color=COLOR_TEXT)
    
    return slide

# ==================== 绗?椤碉細鐗╃悊寮曟搸 ====================

def create_slide_5_physics(prs):
    """鏍稿績鎶€鏈?鈥斺€旂墿鐞嗗紩鎿?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "鐗╃悊寮曟搸锛氬弻鍏辨尟妯″瀷绮惧噯鎷熷悎 Fano 鍏辨尟",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 宸︽爮锛氭礇浼﹀吂 ED+MD 鍙屽叡鎸ā鍨?
    left_x = Inches(0.8)
    col_width = Inches(5.8)
    
    add_title_shape(slide, "娲涗鸡鍏?ED+MD 鍙屽叡鎸ā鍨?,
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 鐞嗚鍩虹
    theory_text = ("鐞嗚鍩虹锛氬熀浜庣背姘忔暎灏勭悊璁猴紝浜氭尝闀跨撼绫虫煴鍙悓鏃舵縺鍙?
                  "**鐢靛伓鏋佸瓙 (ED)** 涓?**纾佸伓鏋佸瓙 (MD)** 涓ょ鍏辨尟妯″紡")
    add_body_text(slide, theory_text,
                  left_x, Inches(2.0), col_width, Inches(0.8),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 鏍稿績鍘熺悊
    principle_text = ("鏍稿績鍘熺悊锛欵D 涓?MD 鐨勭浉骞插彔鍔犲彲浜х敓闈炲绉扮殑 Fano 鍏辨尟绾垮瀷锛?
                     "绯荤粺閫氳繃鍙屾礇浼﹀吂鎸瓙妯″瀷鍒嗗埆鎷熷悎涓ょ鍏辨尟锛屽疄鐜板弽灏勫厜璋辩殑蹇€熻В鏋愯绠?)
    add_body_text(slide, principle_text,
                  left_x, Inches(2.8), col_width, Inches(0.8),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 鏍稿績鍏紡
    add_body_text(slide, "鏍稿績鍏紡锛?,
                  left_x, Inches(3.7), col_width, Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    
    # 鍏紡锛堢畝鍖栬〃绀猴紝瀹為檯搴斾娇鐢ㄥ叕寮忕紪杈戝櫒锛?
    formula = "蟽_sca = (2蟺/k虏) 危_{l=1}^{鈭瀩 (2l+1)(|a_l|虏 + |b_l|虏)"
    formula_note = "鍏朵腑 a_l 涓虹數澶氭瀬瀛愭暎灏勭郴鏁帮紝b_l 涓虹澶氭瀬瀛愭暎灏勭郴鏁帮紝k 涓鸿嚜鐢辩┖闂存尝鏁?
    
    add_body_text(slide, formula,
                  left_x + Inches(0.3), Inches(4.1), col_width - Inches(0.3), Inches(0.5),
                  font_size=Pt(14), color=COLOR_TEXT)
    
    add_body_text(slide, formula_note,
                  left_x + Inches(0.3), Inches(4.6), col_width - Inches(0.3), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT)
    
    # 浼樺娍
    advantage = "浼樺娍锛氱浉姣斿叏娉豢鐪燂紝璁＄畻閫熷害鎻愬崌 3 涓暟閲忕骇锛岀簿搴︽弧瓒冲伐绋嬭璁￠渶姹?
    add_body_text(slide, advantage,
                  left_x, Inches(5.3), col_width, Inches(0.5),
                  font_size=Pt(12), bold=True, color=COLOR_ACCENT)
    
    # 宸︽爮閰嶅浘鍗犱綅
    add_rectangle(slide, left_x, Inches(5.9), Inches(3.0), Inches(1.2),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[ED/MD 鍏辨尟鍦哄己鍒嗗竷绀烘剰鍥綸",
                  left_x, Inches(7.15), Inches(3.0), Inches(0.3),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 鍙虫爮锛欶P 鑵旇壊鍩熸墿灞?
    right_x = Inches(6.8)
    add_title_shape(slide, "FP 鑵旇壊鍩熸墿灞曟ā鍧?,
                    right_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 缁撴瀯
    structure = "缁撴瀯锛欰g (30 nm) / TiO鈧?T) / Ag (bulk) 閲戝睘-浠嬭川-閲戝睘涓夋槑娌荤粨鏋?
    add_body_text(slide, structure,
                  right_x, Inches(2.0), col_width, Inches(0.5),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 鍘熺悊
    fp_principle = ("鍘熺悊锛氬熀浜庡鍏夋潫骞叉秹鍘熺悊锛岄€氳繃璋冭妭浠嬭川鑵斿帤搴?T 璋冩帶鍏辨尟娉㈤暱锛?
                   "瀹炵幇瀹藉厜璋辫寖鍥寸殑棰滆壊璋冩帶")
    add_body_text(slide, fp_principle,
                  right_x, Inches(2.6), col_width, Inches(0.7),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 浠峰€?
    value = ("浠峰€硷細绐佺牬鍗曚竴 TiO鈧?绾崇背鏌辩殑鑹插煙鐗╃悊闄愬埗锛?
             "鍙疄鐜伴珮楗卞拰搴﹂潚钃濊壊涓庣函绾㈣壊")
    add_body_text(slide, value,
                  right_x, Inches(3.4), col_width, Inches(0.6),
                  font_size=Pt(12), bold=True, color=COLOR_ACCENT)
    
    # 鍙虫爮閰嶅浘鍗犱綅
    add_rectangle(slide, right_x, Inches(4.2), Inches(3.5), Inches(2.5),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[FP 鑵旂粨鏋勫墫闈㈠浘]",
                  right_x, Inches(6.75), Inches(3.5), Inches(0.3),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 绗?椤碉細ML鍔犻€熷紩鎿?====================

def create_slide_6_ml(prs):
    """鏍稿績鎶€鏈?鈥斺€擬L鍔犻€熷紩鎿?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "ML 鍔犻€燂細ResMLP 娈嬪樊浠ｇ悊妯″瀷瀹炵幇姣绾ч娴?,
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 宸︽爮锛氭ā鍨嬫灦鏋?
    left_x = Inches(0.8)
    col_width = Inches(5.8)
    
    add_title_shape(slide, "妯″瀷鏋舵瀯",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    arch_items = [
        "缃戠粶缁撴瀯锛氶噰鐢?ResMLP 娈嬪樊澶氬眰鎰熺煡鏈猴紝鍖呭惈 4 涓?256 缁存畫宸潡锛?
        "閫氳繃娈嬪樊杩炴帴閬垮厤姊害娑堝け锛屼繚璇佹繁灞傜綉缁滆缁冪ǔ瀹氭€?,
        "杈撳叆缁村害锛? 缁寸壒寰?鈥?绾崇背鏌辩洿寰?D銆侀珮搴?H銆佸懆鏈?P銆佹潗鏂欐爣璇嗐€佹尝闀跨储寮?,
        "杈撳嚭缁村害锛?1 涓尝闀跨偣鐨勫弽灏勭巼鍊硷紝瑕嗙洊 400-800 nm 鍙鍏夋尝娈碉紝姝ラ暱 5 nm",
        "閮ㄧ讲鏂瑰紡锛氭ā鍨嬪鍑轰负 ONNX 鏍煎紡锛孋PU 鍗冲彲瀹屾垚鏋侀€熸帹鐞嗭紝鏃犻渶 GPU 鏀寔"
    ]
    
    y_pos = Inches(2.0)
    for item in arch_items:
        # 椤圭洰绗﹀彿
        add_rectangle(slide, left_x + Inches(0.1), y_pos + Inches(0.08), 
                     Inches(0.08), Inches(0.08),
                     fill_color=COLOR_PRIMARY, line_color=None)
        add_body_text(slide, item,
                     left_x + Inches(0.3), y_pos, col_width - Inches(0.4), Inches(0.9),
                     font_size=Pt(12), color=COLOR_TEXT)
        y_pos += Inches(0.85)
    
    # 宸︽爮閰嶅浘
    add_rectangle(slide, left_x, Inches(5.5), Inches(4.5), Inches(1.6),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[ResMLP 缃戠粶缁撴瀯绀烘剰鍥綸",
                  left_x + Inches(0.75), Inches(7.15), Inches(3.0), Inches(0.3),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 鍙虫爮锛氭€ц兘鎸囨爣
    right_x = Inches(6.8)
    
    add_title_shape(slide, "鎬ц兘鎸囨爣",
                    right_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 鏍稿績鎸囨爣鍗＄墖
    metrics = [
        ("鎺ㄧ悊閫熷害", "< 1 ms / 鏉″厜璋?, 
         "鐩告瘮娲涗鸡鍏硅В鏋愭ā鍨嬪啀鎻愬崌 2 涓暟閲忕骇锛岀湡姝ｅ疄鐜板疄鏃跺搷搴?),
        ("棰勬祴绮惧害", "螖E鈧傗個鈧€鈧€ < 2", 
         "鑹插樊浣庝簬浜虹溂鍙垎杈ㄩ槇鍊硷紝杈惧埌宸ヤ笟绾у簲鐢ㄧ簿搴?),
        ("璁粌鏁版嵁闆?, "鍗佷竾绾ф爣娉ㄦ牱鏈?, 
         "鍩轰簬鐗╃悊寮曟搸鐢熸垚锛屼繚璇佺墿鐞嗕竴鑷存€?)
    ]
    
    y_pos = Inches(2.0)
    for title, value, desc in metrics:
        # 鎸囨爣鍗＄墖鑳屾櫙
        add_rectangle(slide, right_x, y_pos, col_width, Inches(1.3),
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(1.5))
        
        # 鏍囬
        add_body_text(slide, title,
                     right_x + Inches(0.2), y_pos + Inches(0.15), 
                     Inches(1.5), Inches(0.35),
                     font_size=Pt(12), bold=True, color=COLOR_TEXT)
        
        # 鏁板€硷紙鐞ョ弨閲戦珮浜斁澶э級
        add_title_shape(slide, value,
                       right_x + Inches(0.2), y_pos + Inches(0.5), 
                       col_width - Inches(0.4), Inches(0.5),
                       font_size=Pt(20), bold=True, color=COLOR_ACCENT)
        
        # 璇存槑
        add_body_text(slide, desc,
                     right_x + Inches(0.2), y_pos + Inches(1.0), 
                     col_width - Inches(0.4), Inches(0.3),
                     font_size=Pt(10), color=COLOR_TEXT)
        
        y_pos += Inches(1.5)
    
    # 棰勬祴鍏夎氨瀵规瘮鍥?
    add_rectangle(slide, right_x, Inches(6.0), Inches(5.5), Inches(1.1),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[棰勬祴鍏夎氨 vs 浠跨湡鍏夎氨 閲嶅悎瀵规瘮鏇茬嚎]",
                  right_x + Inches(1.5), Inches(7.15), Inches(2.5), Inches(0.3),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 绗?椤碉細澶氱畻娉曢€嗚璁?====================

def create_slide_7_inverse(prs):
    """鏍稿績鎶€鏈?鈥斺€斿绠楁硶閫嗚璁′綋绯?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "閫嗚璁★細鍥涚绠楁硶瑕嗙洊绮惧害-閫熷害鍏ㄥ満鏅?,
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 涓婂崐娈靛畾涔?
    definition = ("閫嗚璁″嵆 \"浠ョ洰鏍囬鑹插弽鎺ㄧ粨鏋勫弬鏁癨"锛?
                 "杈撳叆鐩爣 RGB / 鍗佸叚杩涘埗鑹插€硷紝鑷姩姹傝В鏈€浼樼殑绾崇背鏌辩洿寰勩€侀珮搴︺€佸懆鏈熺粍鍚堬紝"
                 "瀹炵幇 \"鎵€鎯冲嵆鎵€寰梊" 鐨勮璁′綋楠屻€?)
    add_body_text(slide, definition,
                  Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.7),
                  font_size=Pt(14), color=COLOR_TEXT)
    
    # 鍥涘垪绠楁硶瀵规瘮鍗＄墖
    card_width = Inches(2.9)
    card_height = Inches(3.2)
    card_y = Inches(2.3)
    
    algorithms = [
        ("缃戞牸鎼滅储娉?, 
         "鍏ㄥ弬鏁扮┖闂寸┓涓鹃亶鍘哱n锛堢害 11000 缁勫弬鏁帮級",
         "绾?30 s", "螖E鈧傗個鈧€鈧€ 鈮?2.1", "灏忚寖鍥寸簿鍑嗗浼樸€佸熀鍑嗛獙璇?),
        ("姊害涓嬮檷娉?,
         "鍩轰簬浠ｇ悊妯″瀷鐨刓nAdam 姊害浼樺寲",
         "绾?12 s", "螖E鈧傗個鈧€鈧€ 鈮?3.8", "鍗曠洰鏍囧揩閫熶紭鍖?),
        ("寮哄寲瀛︿範娉?,
         "Q-Learning 鏅鸿兘浣揬n绂绘暎鍐崇瓥鎼滅储",
         "绾?5 s", "螖E鈧傗個鈧€鈧€ 鈮?5.7", "澶ц寖鍥村弬鏁扮┖闂村垵姝ョ瓫閫?),
        ("浠ｇ悊妯″瀷娉?,
         "ResMLP 缃戠粶鐩存帴\n閫嗗悜鎺ㄧ悊鏈€浼樺弬鏁?,
         "< 1 s", "螖E鈧傗個鈧€鈧€ 鈮?3.8", "缁濆ぇ澶氭暟甯歌璁捐鍦烘櫙")
    ]
    
    headers = ["鍘熺悊", "鑰楁椂", "绮惧害", "閫傜敤鍦烘櫙"]
    
    for i, (name, principle, time, accuracy, scenario) in enumerate(algorithms):
        x_pos = Inches(0.8) + i * (card_width + Inches(0.25))
        
        # 鍗＄墖鑳屾櫙
        add_rectangle(slide, x_pos, card_y, card_width, card_height,
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(2))
        
        # 绠楁硶鍚嶇О锛堥《閮ㄨ壊鍧楋級
        add_rectangle(slide, x_pos, card_y, card_width, Inches(0.5),
                     fill_color=COLOR_PRIMARY, line_color=None)
        add_title_shape(slide, name,
                       x_pos, card_y + Inches(0.1), card_width, Inches(0.4),
                       font_size=Pt(15), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                       align=PP_ALIGN.CENTER)
        
        # 鍐呭琛?
        rows = [
            ("鍘熺悊", principle),
            ("鑰楁椂", time),
            ("绮惧害", accuracy),
            ("閫傜敤鍦烘櫙", scenario)
        ]
        
        row_y = card_y + Inches(0.6)
        for label, content in rows:
            # 鏍囩
            add_body_text(slide, label + "锛?,
                         x_pos + Inches(0.15), row_y, Inches(0.8), Inches(0.25),
                         font_size=Pt(10), bold=True, color=COLOR_PRIMARY)
            # 鍐呭
            content_color = COLOR_ACCENT if label in ["鑰楁椂", "绮惧害"] else COLOR_TEXT
            is_bold = label in ["鑰楁椂", "绮惧害"]
            add_body_text(slide, content,
                         x_pos + Inches(0.15), row_y + Inches(0.25), 
                         card_width - Inches(0.3), Inches(0.55),
                         font_size=Pt(11), bold=is_bold, color=content_color)
            row_y += Inches(0.85)
    
    # 搴曢儴閫熷害-绮惧害鏁ｇ偣鍥?
    scatter_y = Inches(5.8)
    add_rectangle(slide, Inches(3.0), scatter_y, Inches(7.0), Inches(1.4),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    
    # 鍧愭爣杞存爣娉?
    add_body_text(slide, "閫熷害 鈫?,
                  Inches(6.0), Inches(7.15), Inches(1.0), Inches(0.3),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    add_body_text(slide, "绮惧害 鈫?,
                  Inches(2.5), Inches(6.3), Inches(0.5), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT)
    
    add_body_text(slide, "[閫熷害-绮惧害浜岀淮鏁ｇ偣鍥撅細鍥涚绠楁硶瀹氫綅]",
                  Inches(4.5), Inches(7.15), Inches(4.0), Inches(0.3),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 绗?椤碉細鏍稿績鍔熻兘浣撶郴锛堜笂锛?====================

def create_slide_8_functions_1(prs):
    """鏍稿績鍔熻兘浣撶郴锛堜笂锛?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "鏍稿績鍔熻兘锛氬疄鏃堕瑙堜笌鏅鸿兘閫嗚璁?,
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 宸︽爮锛氬崟鏌?鍙屾煴瀹炴椂棰勮
    left_x = Inches(0.8)
    col_width = Inches(5.8)
    
    add_title_shape(slide, "鍗曟煴 / 鍙屾煴瀹炴椂棰勮",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 鎿嶄綔鏂瑰紡
    add_body_text(slide, "鎿嶄綔鏂瑰紡锛?,
                  left_x, Inches(2.0), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "婊戝姩鏉¤皟鑺傜洿寰?D銆侀珮搴?H銆佸懆鏈?P 鍙傛暟锛屾绉掔骇鍚屾鏇存柊缁撴灉",
                  left_x + Inches(1.2), Inches(2.0), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 杈撳嚭鍐呭
    add_body_text(slide, "杈撳嚭鍐呭锛?,
                  left_x, Inches(2.7), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "瀹炴椂鏄剧ず鍙嶅皠鑹插崱銆佸畬鏁村弽灏勫厜璋辨洸绾裤€佸叧閿叡鎸尝闀挎爣娉?,
                  left_x + Inches(1.2), Inches(2.7), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 鎵╁睍鑳藉姏
    add_body_text(slide, "鎵╁睍鑳藉姏锛?,
                  left_x, Inches(3.4), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "鏀寔 TE/TM 鍋忔尟鍒囨崲锛屾敮鎸?TiO鈧傘€乤-Si銆丄g銆丄l 鍥涚鏉愭枡涓€閿垏鎹㈠姣?,
                  left_x + Inches(1.2), Inches(3.4), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 鐣岄潰鎴浘鍗犱綅
    add_rectangle(slide, left_x, Inches(4.2), col_width, Inches(3.0),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[绯荤粺鍔熻兘鐣岄潰鎴浘锛氬疄鏃堕瑙堟ā鍧梋\n鏍囨敞鏍稿績鎿嶄綔鍖哄煙",
                  left_x + Inches(1.5), Inches(5.5), Inches(2.8), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 鍙虫爮锛氱綉鏍兼悳绱㈤€嗚璁?
    right_x = Inches(6.8)
    
    add_title_shape(slide, "缃戞牸鎼滅储閫嗚璁?,
                    right_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 鎿嶄綔鏂瑰紡
    add_body_text(slide, "鎿嶄綔鏂瑰紡锛?,
                  right_x, Inches(2.0), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "杈撳叆鐩爣鍗佸叚杩涘埗鑹插€兼垨 RGB 鍊硷紝涓€閿惎鍔ㄨ嚜鍔ㄦ悳绱?,
                  right_x + Inches(1.2), Inches(2.0), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 杈撳嚭缁撴灉
    add_body_text(slide, "杈撳嚭缁撴灉锛?,
                  right_x, Inches(2.7), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "杩斿洖 Top3 鏈€浼樼粨鏋勫弬鏁帮紝鏍囨敞瀵瑰簲鑹插樊 螖E鈧傗個鈧€鈧€",
                  right_x + Inches(1.2), Inches(2.7), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 鏅鸿兘鎻愮ず
    add_body_text(slide, "鏅鸿兘鎻愮ず锛?,
                  right_x, Inches(3.4), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "鑷姩鏍囨敞褰撳墠鏉愭枡鐨勮壊鍩熻竟鐣岋紝缁欏嚭鐗╃悊鍙鎬ф彁绀轰笌鏉愭枡鍒囨崲寤鸿",
                  right_x + Inches(1.2), Inches(3.4), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 鐣岄潰鎴浘鍗犱綅
    add_rectangle(slide, right_x, Inches(4.2), col_width, Inches(3.0),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[绯荤粺鍔熻兘鐣岄潰鎴浘锛氶€嗚璁℃ā鍧梋\n鏍囨敞鏍稿績鎿嶄綔鍖哄煙",
                  right_x + Inches(1.5), Inches(5.5), Inches(2.8), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 绗?椤碉細鏍稿績鍔熻兘浣撶郴锛堜笅锛?====================

def create_slide_9_functions_2(prs):
    """鏍稿績鍔熻兘浣撶郴锛堜笅锛?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "鏍稿績鍔熻兘锛氶珮绾у垎鏋愪笌鎵归噺璁捐",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 涓夊垪妯悜鍗＄墖
    card_width = Inches(3.9)
    card_height = Inches(5.5)
    card_y = Inches(1.5)
    
    functions = [
        ("瑙掕氨杩滃満浼犳挱",
         "妯℃嫙 0掳-60掳 瑙傛祴瑙掑害涓嬬殑鍏夎氨涓庨鑹插彉鍖?,
         "鍒嗘瀽缁撴瀯鑹茬殑瑙掑害渚濊禆鎬э紝璇勪及澶ц瑙掍笅鐨勮壊褰╃ǔ瀹氭€?),
        ("宸ヨ壓瀹瑰樊鍒嗘瀽",
         "鑷姩璁＄畻 卤5 nm 鍙傛暟鍋忓樊涓嬬殑鑹插樊鍙樺寲",
         "璇勪及璁捐鐨勫伐鑹洪瞾妫掓€э紝鎸囧瀹為檯寰撼鍔犲伐鍒跺锛岄檷浣庡埗澶囧け璐ラ闄?),
        ("鍥炬鎵归噺鐢熸垚",
         "瀵煎叆鑷畾涔夊儚绱犲浘妗堬紝鎵归噺鐢熸垚瀵瑰簲缁撴瀯鍙傛暟闃靛垪",
         "鐩存帴瀵规帴寰撼鍔犲伐娴佺▼锛屽ぇ骞呯缉鐭粠璁捐鍒板埗澶囩殑閾捐矾")
    ]
    
    for i, (name, func, value) in enumerate(functions):
        x_pos = Inches(0.8) + i * (card_width + Inches(0.3))
        
        # 鍗＄墖鑳屾櫙
        add_rectangle(slide, x_pos, card_y, card_width, card_height,
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(2))
        
        # 鍔熻兘鍚嶇О锛堥《閮ㄨ壊鍧楋級
        add_rectangle(slide, x_pos, card_y, card_width, Inches(0.6),
                     fill_color=COLOR_PRIMARY, line_color=None)
        add_title_shape(slide, name,
                       x_pos, card_y + Inches(0.15), card_width, Inches(0.4),
                       font_size=Pt(16), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                       align=PP_ALIGN.CENTER)
        
        # 鍔熻兘鎻忚堪
        add_body_text(slide, "鍔熻兘锛?,
                     x_pos + Inches(0.2), card_y + Inches(0.8), 
                     Inches(0.6), Inches(0.35),
                     font_size=Pt(11), bold=True, color=COLOR_PRIMARY)
        add_body_text(slide, func,
                     x_pos + Inches(0.2), card_y + Inches(1.15), 
                     card_width - Inches(0.4), Inches(0.9),
                     font_size=Pt(12), color=COLOR_TEXT)
        
        # 浠峰€?
        add_body_text(slide, "浠峰€硷細",
                     x_pos + Inches(0.2), card_y + Inches(2.1), 
                     Inches(0.6), Inches(0.35),
                     font_size=Pt(11), bold=True, color=COLOR_ACCENT)
        add_body_text(slide, value,
                     x_pos + Inches(0.2), card_y + Inches(2.45), 
                     card_width - Inches(0.4), Inches(1.0),
                     font_size=Pt(12), color=COLOR_TEXT)
        
        # 绀烘剰鍥惧崰浣?
        add_rectangle(slide, x_pos + Inches(0.3), card_y + Inches(3.6), 
                     card_width - Inches(0.6), Inches(1.6),
                     fill_color=RGBColor(0xE8, 0xEA, 0xED),
                     line_color=COLOR_PRIMARY, line_width=Pt(1))
        add_body_text(slide, f"[{name}绀烘剰鍥綸",
                     x_pos + Inches(0.5), card_y + Inches(5.25), 
                     card_width - Inches(1.0), Inches(0.3),
                     font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 绗?0椤碉細鏉愭枡涓庤‖搴曚綋绯?====================

# ==================== 绗?0椤碉細鏉愭枡涓庤‖搴曚綋绯?====================

def create_slide_10_materials(prs):
    """鏉愭枡涓庤‖搴曚綋绯?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "澶氭潗鏂欎綋绯婚€傞厤涓嶅悓搴旂敤鍦烘櫙",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 宸︽爮锛氭潗鏂欑壒鎬у姣旇〃
    left_x = Inches(0.8)
    table_width = Inches(5.5)
    
    add_title_shape(slide, "鏉愭枡鐗规€у姣?,
                    left_x, Inches(1.4), table_width, Inches(0.5),
                    font_size=Pt(16), bold=True, color=COLOR_PRIMARY)
    
    # 琛ㄦ牸鏁版嵁 - 绮剧‘鍖归厤鐢ㄦ埛鎻愪緵鐨勮〃鏍?
    table_data = [
        ["绫诲埆", "鍙€夋潗鏂?, "鏍稿績鍏夊鐗圭偣", "鍏稿瀷搴旂敤鍦烘櫙"],
        ["浠嬭川鏉愭枡", "TiO2锛堥攼閽涚熆锛夈€乤-Si锛堥潪鏅剁锛?, "鍏夊鎹熻€椾綆銆佽壊褰╂煍鍜?, "閫忓皠寮忔樉绀恒€佺敓鐗╀紶鎰?],
        ["閲戝睘鏉愭枡", "Ag銆丄l", "鍏辨尟寮哄害楂樸€佽壊褰╅ケ鍜屽害楂?, "鍙嶅皠寮忔樉绀恒€佸厜瀛﹂槻浼?],
        ["琛簳鏉愭枡", "SiO2銆丼i3N4", "宸ヨ壓鎴愮啛銆佸櫒浠跺吋瀹规€уソ", "鏍囧噯寰撼鍔犲伐宸ヨ壓"]
    ]
    
    # 琛ㄥご
    header_height = Inches(0.45)
    row_height = Inches(0.55)
    col_widths = [Inches(1.0), Inches(1.6), Inches(1.5), Inches(1.4)]
    
    y_pos = Inches(2.0)
    
    # 琛ㄥご鑳屾櫙
    add_rectangle(slide, left_x, y_pos, table_width, header_height,
                 fill_color=COLOR_PRIMARY, line_color=None)
    
    x_offset = left_x
    for j, cell in enumerate(table_data[0]):
        add_body_text(slide, cell,
                     x_offset + Inches(0.05), y_pos + Inches(0.1), 
                     col_widths[j] - Inches(0.1), header_height - Inches(0.2),
                     font_size=Pt(10), bold=True, 
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        x_offset += col_widths[j]
    
    y_pos += header_height
    
    # 琛ㄦ牸鍐呭
    for i, row in enumerate(table_data[1:]):
        # 浜ゆ浛琛岃儗鏅壊
        bg_color = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xF2, 0xF3, 0xF5)
        add_rectangle(slide, left_x, y_pos, table_width, row_height,
                     fill_color=bg_color, line_color=COLOR_PRIMARY, line_width=Pt(0.5))
        
        x_offset = left_x
        for j, cell in enumerate(row):
            # 绗竴鍒楀姞绮?
            is_bold = (j == 0)
            font_sz = Pt(9) if j == 1 else Pt(10)  # 鏉愭枡鍒楀瓧浣撶◢灏忎互瀹圭撼鍐呭
            
            add_body_text(slide, cell,
                         x_offset + Inches(0.05), y_pos + Inches(0.12), 
                         col_widths[j] - Inches(0.1), row_height - Inches(0.24),
                         font_size=font_sz, bold=is_bold, color=COLOR_TEXT, align=PP_ALIGN.LEFT)
            x_offset += col_widths[j]
        
        y_pos += row_height
    
    # 鍙虫爮锛氳壊鍩熷姣?
    right_x = Inches(6.8)
    right_width = Inches(5.8)
    
    add_title_shape(slide, "鑹插煙瀵规瘮鍒嗘瀽",
                    right_x, Inches(1.4), right_width, Inches(0.5),
                    font_size=Pt(16), bold=True, color=COLOR_PRIMARY)
    
    # 璇存槑鏂囧瓧
    desc_text = ("鐩磋灞曠ず绾?TiO2 绾崇背鏌变笌 FP 鑵旂粨鏋勭殑鑹插煙瑕嗙洊鑼冨洿宸紓\n"
                "鏍囨敞姣忕鏉愭枡鐨勭墿鐞嗚壊鍩熸瀬闄愶紝浣撶幇绯荤粺鐨勭墿鐞嗕弗璋ㄦ€?)
    add_body_text(slide, desc_text,
                  right_x, Inches(2.0), right_width, Inches(0.8),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 鏍稿績缁撹锛堥珮浜級
    conclusion = "鏍稿績缁撹锛欶P 鑵旂粨鏋勫彲灏嗚壊鍩熻鐩栬寖鍥存彁鍗囩害 40%"
    add_title_shape(slide, conclusion,
                    right_x, Inches(2.9), right_width, Inches(0.5),
                    font_size=Pt(14), bold=True, color=COLOR_ACCENT)
    
    # CIE 1931 鑹插搧鍥惧崰浣?
    add_rectangle(slide, right_x, Inches(3.5), right_width, Inches(3.5),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[CIE 1931 鑹插搧鍥綸\n鏍囨敞 TiO2 绾崇背鏌变笌 FP 鑵旂粨鏋勭殑鑹插煙杈圭晫",
                  right_x + Inches(1.2), Inches(5.0), Inches(3.4), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 鍥句緥璇存槑
    legend_y = Inches(7.0)
    add_rectangle(slide, right_x + Inches(0.5), legend_y, Inches(0.4), Inches(0.2),
                 fill_color=COLOR_PRIMARY, line_color=None)
    add_body_text(slide, "TiO2 绾崇背鏌?,
                  right_x + Inches(1.0), legend_y, Inches(1.5), Inches(0.25),
                  font_size=Pt(10), color=COLOR_TEXT)
    
    add_rectangle(slide, right_x + Inches(2.8), legend_y, Inches(0.4), Inches(0.2),
                 fill_color=COLOR_ACCENT, line_color=None)
    add_body_text(slide, "FP 鑵旂粨鏋?,
                  right_x + Inches(3.3), legend_y, Inches(1.5), Inches(0.25),
                  font_size=Pt(10), color=COLOR_TEXT)
    
    return slide

# ==================== 绗?1椤碉細鎬ц兘楠岃瘉 ====================

def create_slide_11_performance(prs):
    """鎬ц兘楠岃瘉涓庝紭鍔垮姣?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "鎬ц兘鎸囨爣涓庢柟妗堝姣?,
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 宸︽爮锛氭牳蹇冮噺鍖栨寚鏍?
    left_x = Inches(0.8)
    col_width = Inches(5.5)
    
    add_title_shape(slide, "鏍稿績閲忓寲鎸囨爣",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    metrics = [
        ("鍗曟潯鍏夎氨棰勬祴閫熷害", "< 1 ms"),
        ("浠ｇ悊妯″瀷閫嗚璁＄簿搴?, "DeltaE2000 = 3.8"),
        ("鏀寔鏉愭枡+琛簳缁勫悎", "4 + 3 绉?),
        ("瀹㈡埛绔畨瑁呬緷璧?, "0 涓?)
    ]
    
    y_pos = Inches(2.1)
    for label, value in metrics:
        # 鎸囨爣鍗＄墖
        add_rectangle(slide, left_x, y_pos, col_width, Inches(0.9),
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(1.5))
        
        # 鏍囩
        add_body_text(slide, label,
                     left_x + Inches(0.2), y_pos + Inches(0.15), 
                     Inches(2.5), Inches(0.35),
                     font_size=Pt(12), color=COLOR_TEXT)
        
        # 鏁板€硷紙鐞ョ弨閲戦珮浜斁澶э級
        add_title_shape(slide, value,
                       left_x + Inches(2.8), y_pos + Inches(0.2), 
                       Inches(2.5), Inches(0.55),
                       font_size=Pt(22), bold=True, color=COLOR_ACCENT)
        
        y_pos += Inches(1.05)
    
    # 鍙虫爮锛氫笌浼犵粺璁捐鏂规瀵规瘮
    right_x = Inches(6.8)
    right_width = Inches(5.8)
    
    add_title_shape(slide, "涓庝紶缁熻璁℃柟妗堝姣?,
                    right_x, Inches(1.4), right_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 瀵规瘮琛ㄦ暟鎹?- 绮剧‘鍖归厤鐢ㄦ埛鎻愪緵鐨勮〃鏍?
    compare_data = [
        ["瀵规瘮缁村害", "浼犵粺 FDTD 浠跨湡", "鏈郴缁?],
        ["鍗曠粨鏋勮绠楁椂闂?, "鏁板皬鏃?, "< 1 ms"],
        ["杞欢渚濊禆", "鍟嗕笟鎺堟潈杞欢", "鏃?],
        ["瀛︿範闂ㄦ", "鏁版湀涓撲笟璁粌", "鍗冲紑鍗崇敤"],
        ["閫嗚璁¤兘鍔?, "鏃?, "鍥涚绠楁硶"],
        ["璁惧瑕佹眰", "楂樻€ц兘宸ヤ綔绔?, "鏅€氭祻瑙堝櫒"]
    ]
    
    header_height = Inches(0.4)
    row_height = Inches(0.52)
    col_widths = [Inches(1.7), Inches(1.9), Inches(2.2)]
    
    y_pos = Inches(2.1)
    
    # 琛ㄥご
    add_rectangle(slide, right_x, y_pos, right_width, header_height,
                 fill_color=COLOR_PRIMARY, line_color=None)
    x_offset = right_x
    for j, cell in enumerate(compare_data[0]):
        add_body_text(slide, cell,
                     x_offset + Inches(0.05), y_pos + Inches(0.08), 
                     col_widths[j] - Inches(0.1), header_height - Inches(0.16),
                     font_size=Pt(10), bold=True, 
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        x_offset += col_widths[j]
    
    y_pos += header_height
    
    # 琛ㄦ牸鍐呭
    for i, row in enumerate(compare_data[1:]):
        bg_color = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xF2, 0xF3, 0xF5)
        add_rectangle(slide, right_x, y_pos, right_width, row_height,
                     fill_color=bg_color, line_color=COLOR_PRIMARY, line_width=Pt(0.5))
        
        x_offset = right_x
        for j, cell in enumerate(row):
            # 鏈郴缁熷垪楂樹寒
            text_color = COLOR_ACCENT if j == 2 else COLOR_TEXT
            is_bold = (j == 2)
            add_body_text(slide, cell,
                         x_offset + Inches(0.05), y_pos + Inches(0.1), 
                         col_widths[j] - Inches(0.1), row_height - Inches(0.2),
                         font_size=Pt(10), bold=is_bold, 
                         color=text_color, align=PP_ALIGN.CENTER)
            x_offset += col_widths[j]
        
        y_pos += row_height
    
    # 搴曢儴鎬荤粨
    summary = "鍦ㄤ繚璇佸伐绋嬪彲鐢ㄧ簿搴︾殑鍓嶆彁涓嬶紝瀹炵幇璁捐鏁堢巼鐨勬寚鏁扮骇鎻愬崌"
    add_title_shape(slide, summary,
                    Inches(2.5), Inches(6.8), Inches(8.0), Inches(0.5),
                    font_size=Pt(16), bold=True, color=COLOR_PRIMARY)
    
    return slide
def create_slide_12_cloud(prs):
    """浜戠閮ㄧ讲涓庝娇鐢ㄦ柟寮?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "浜戠閮ㄧ讲锛氶浂闂ㄦ寮€绠卞嵆鐢?,
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 宸︽爮锛氬畬鏁存妧鏈爤
    left_x = Inches(0.8)
    col_width = Inches(4.5)
    
    add_title_shape(slide, "瀹屾暣鎶€鏈爤",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    tech_stack = [
        ("鍓嶇妗嗘灦", "Streamlit Python Web"),
        ("鐗╃悊寮曟搸", "娲涗鸡鍏?ED+MD 鍙屽叡鎸ā鍨?),
        ("ML 寮曟搸", "ResMLP (PyTorch) + ONNX"),
        ("澶фā鍨?, "DeepSeek API"),
        ("閮ㄧ讲骞冲彴", "GitHub + Streamlit Cloud")
    ]
    
    y_pos = Inches(2.0)
    for label, value in tech_stack:
        add_body_text(slide, f"{label}锛?,
                     left_x, y_pos, Inches(1.3), Inches(0.4),
                     font_size=Pt(12), bold=True, color=COLOR_PRIMARY)
        add_body_text(slide, value,
                     left_x + Inches(1.4), y_pos, Inches(3.0), Inches(0.4),
                     font_size=Pt(12), color=COLOR_TEXT)
        y_pos += Inches(0.6)
    
    # 鍙虫爮锛氫笁姝ヤ娇鐢ㄦ祦绋?
    right_x = Inches(8.0)
    
    add_title_shape(slide, "涓夋浣跨敤娴佺▼",
                    right_x, Inches(1.4), Inches(4.5), Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    steps = [
        ("鈶?, "娴忚鍣ㄦ墦寮€浜戠鍦板潃锛屾棤闇€涓嬭浇瀹夎浠讳綍杞欢"),
        ("鈶?, "璋冭妭鍙傛暟鎴栬緭鍏ョ洰鏍囬鑹诧紝瀹炴椂鑾峰彇鍏夎氨涓庨鑹茬粨鏋?),
        ("鈶?, "涓€閿皟鐢ㄥぇ妯″瀷鍒嗘瀽锛岃幏鍙栫墿鐞嗚В閲婁笌浼樺寲寤鸿")
    ]
    
    y_pos = Inches(2.0)
    for num, desc in steps:
        # 搴忓彿鍦嗗湀
        add_rectangle(slide, right_x, y_pos, Inches(0.5), Inches(0.5),
                     fill_color=COLOR_ACCENT, line_color=None)
        add_title_shape(slide, num,
                       right_x, y_pos + Inches(0.05), Inches(0.5), Inches(0.45),
                       font_size=Pt(16), bold=True, 
                       color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        
        # 鎻忚堪
        add_body_text(slide, desc,
                     right_x + Inches(0.7), y_pos + Inches(0.08), 
                     Inches(3.8), Inches(0.5),
                     font_size=Pt(12), color=COLOR_TEXT)
        
        y_pos += Inches(0.9)
    
    # 涓儴閱掔洰灞曠ず锛氱郴缁熷叆鍙?
    center_y = Inches(5.0)
    add_rectangle(slide, Inches(2.5), center_y, Inches(8.0), Inches(1.6),
                 fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                 line_color=COLOR_ACCENT, line_width=Pt(3))
    
    add_title_shape(slide, "绯荤粺浜戠鍦板潃",
                   Inches(2.5), center_y + Inches(0.15), Inches(8.0), Inches(0.4),
                   font_size=Pt(14), bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
    
    # URL锛堢惀鐝€閲戦珮浜級
    add_title_shape(slide, "https://huggingface.co/spaces/qiaoanqi/metasurface-color-designer",
                   Inches(2.5), center_y + Inches(0.6), Inches(8.0), Inches(0.5),
                   font_size=Pt(16), bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    
    # 浜岀淮鐮佸崰浣?
    add_rectangle(slide, Inches(10.8), center_y + Inches(0.3), 
                 Inches(1.2), Inches(1.2),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[浜岀淮鐮乚",
                  Inches(10.8), center_y + Inches(0.9), Inches(1.2), Inches(0.4),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 绗?3椤碉細鎬荤粨涓庡睍鏈?====================

def create_slide_13_conclusion(prs):
    """鎬荤粨涓庢湭鏉ュ睍鏈?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 椤甸潰澶ф爣棰?
    add_title_shape(slide, "鎬荤粨涓庡睍鏈?,
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(32), bold=True, color=COLOR_PRIMARY)
    
    # 宸︽爮锛氭牳蹇冨垱鏂版€荤粨
    left_x = Inches(0.8)
    col_width = Inches(5.8)
    
    add_title_shape(slide, "鏍稿績鍒涙柊鎬荤粨",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    innovations = [
        "鎻愬嚭 \"鐗╃悊妯″瀷 + 娣卞害瀛︿範\" 鍙岄┍鍔ㄦ灦鏋勶紝鍦ㄧ簿搴︿笌閫熷害闂村疄鐜版渶浼樺钩琛?,
        "鏋勫缓鍥涚閫嗚璁＄畻娉曚綋绯伙紝鍙€傞厤涓嶅悓绮惧害銆侀€熷害闇€姹傜殑璁捐鍦烘櫙",
        "棣栨灏嗗ぇ妯″瀷鍘熺敓寮曞叆瓒呰〃闈㈣璁℃祦绋嬶紝瀹炵幇鏅鸿兘鐗╃悊瑙ｆ瀽涓庝紭鍖栧缓璁?,
        "绾?Web 浜戠閮ㄧ讲妯″紡锛屽ぇ骞呴檷浣庤秴琛ㄩ潰缁撴瀯鑹茬殑璁捐闂ㄦ"
    ]
    
    y_pos = Inches(2.0)
    for i, item in enumerate(innovations, 1):
        # 搴忓彿
        add_rectangle(slide, left_x, y_pos + Inches(0.05), 
                     Inches(0.35), Inches(0.35),
                     fill_color=COLOR_PRIMARY, line_color=None)
        add_title_shape(slide, str(i),
                       left_x, y_pos + Inches(0.08), Inches(0.35), Inches(0.3),
                       font_size=Pt(14), bold=True, 
                       color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        
        # 鍐呭
        add_body_text(slide, item,
                     left_x + Inches(0.5), y_pos, 
                     col_width - Inches(0.6), Inches(0.9),
                     font_size=Pt(13), color=COLOR_TEXT)
        
        y_pos += Inches(1.05)
    
    # 鍙虫爮锛氭湭鏉ュ伐浣滃睍鏈?
    right_x = Inches(6.8)
    
    add_title_shape(slide, "鏈潵宸ヤ綔灞曟湜",
                    right_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    futures = [
        ("鏁版嵁鍗囩骇", "鎺ュ叆 RCWA/FDTD 楂樼簿搴︿豢鐪熸暟鎹缁冿紝杩涗竴姝ユ彁鍗囨ā鍨嬬簿搴?),
        ("缁撴瀯鎵╁睍", "鏀寔 3D 澶嶆潅缁撴瀯銆佸亸鎸皟鎺х粨鏋勭瓑鏇村瓒呰〃闈㈢被鍨?),
        ("璁捐鍗囩骇", "寮曞叆鎵╂暎妯″瀷绛夌敓鎴愬紡 AI锛屽疄鐜版洿鑷敱鐨勫垱鏂扮粨鏋勮璁?),
        ("闂幆鍗囩骇", "瀵规帴瀹為獙鍒跺鏁版嵁锛屾瀯寤?\"璁捐-鍒跺-娴嬭瘯\" 鍏ㄦ祦绋嬮棴鐜?)
    ]
    
    y_pos = Inches(2.0)
    for label, desc in futures:
        # 鏍囩锛堢惀鐝€閲戯級
        add_body_text(slide, f"{label}锛?,
                     right_x, y_pos, Inches(1.2), Inches(0.35),
                     font_size=Pt(12), bold=True, color=COLOR_ACCENT)
        
        # 鎻忚堪
        add_body_text(slide, desc,
                     right_x + Inches(1.3), y_pos, 
                     col_width - Inches(1.4), Inches(0.9),
                     font_size=Pt(12), color=COLOR_TEXT)
        
        y_pos += Inches(1.05)
    
    return slide

# ==================== 绗?4椤碉細鑷磋阿椤?====================

def create_slide_14_thanks(prs):
    """鑷磋阿椤?""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 灞呬腑瀵圭О甯冨眬
    center_x = Inches(6.667)  # 椤甸潰涓績
    
    # 涓绘爣棰?
    add_title_shape(slide, "鎰熻阿鑱嗗惉",
                    Inches(4.0), Inches(2.0), Inches(5.333), Inches(1.0),
                    font_size=Pt(54), bold=True, color=COLOR_PRIMARY, 
                    align=PP_ALIGN.CENTER)
    
    # 鍓爣棰?
    add_body_text(slide, "鏁鍚勪綅鑰佸笀鎵硅瘎鎸囨",
                  Inches(4.0), Inches(3.2), Inches(5.333), Inches(0.6),
                  font_size=Pt(20), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 鍒嗛殧绾?
    add_line(slide, Inches(5.0), Inches(4.0), Inches(8.333), Inches(4.0),
             color=COLOR_PRIMARY, width=Pt(2))
    
    # 鎸囧鏁欏笀
    add_body_text(slide, "鎸囧鏁欏笀锛歑XX 鑰佸笀",
                  Inches(4.0), Inches(4.4), Inches(5.333), Inches(0.5),
                  font_size=Pt(16), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 搴曢儴淇℃伅
    bottom_info = "闀挎矙鐞嗗伐澶у 路 鐗╃悊涓庣數瀛愮瀛﹀闄?路 鍏夌數2501鐝璡n涔斿畨鐞?路 璋㈠鐝?路 渚悽"
    add_body_text(slide, bottom_info,
                  Inches(4.0), Inches(5.2), Inches(5.333), Inches(0.8),
                  font_size=Pt(14), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 鍙充笅瑙掍簯绔湴鍧€
    add_body_text(slide, "绯荤粺璁块棶锛歨ttps://metasurface-web-2you8jsy6wxhmhqgc9vwhm.streamlit.app\n[浜岀淮鐮佸崰浣峕",
                  Inches(8.5), Inches(6.5), Inches(4.0), Inches(0.8),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.RIGHT)
    
    return slide

# ==================== 涓荤▼搴?====================

def main():
    """鐢熸垚瀹屾暣PPT"""
    prs = create_presentation()
    
    # 鍒涘缓14椤靛够鐏墖
    create_slide_1_cover(prs)
    create_slide_2_toc(prs)
    create_slide_3_background(prs)
    create_slide_4_architecture(prs)
    create_slide_5_physics(prs)
    create_slide_6_ml(prs)
    create_slide_7_inverse(prs)
    create_slide_8_functions_1(prs)
    create_slide_9_functions_2(prs)
    create_slide_10_materials(prs)
    create_slide_11_performance(prs)
    create_slide_12_cloud(prs)
    create_slide_13_conclusion(prs)
    create_slide_14_thanks(prs)
    
    # 淇濆瓨鏂囦欢
    output_path = "AI瓒呰〃闈㈢粨鏋勮壊鏅鸿兘璁捐绯荤粺_绛旇京PPT.pptx"
    prs.save(output_path)
    print(f"PPT宸茬敓鎴愶細{os.path.abspath(output_path)}")
    print(f"鍏?{len(prs.slides)} 椤?)

if __name__ == "__main__":
    main()