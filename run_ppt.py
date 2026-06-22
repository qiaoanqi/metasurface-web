"""
AI超表面结构色智能设计系统 - 学术答辩PPT生成器
长沙理工大学 物理与电子科学学院 光电2501班
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# ==================== 全局配置 ====================

# 配色体系
COLOR_PRIMARY = RGBColor(0x16, 0x5D, 0xFF)    # 深空蓝 #165DFF
COLOR_TEXT = RGBColor(0x1D, 0x21, 0x29)        # 炭灰色 #1D2129
COLOR_BG = RGBColor(0xF7, 0xF8, 0xFA)          # 浅灰白 #F7F8FA
COLOR_ACCENT = RGBColor(0xFF, 0x7D, 0x00)      # 琥珀金 #FF7D00

# 字体配置（请确保系统安装思源黑体，或替换为Microsoft YaHei）
FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
FONT_MATH = "Times New Roman"

# 页面尺寸（16:9）
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def create_presentation():
    """创建演示文稿"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # 设置默认背景色
    for layout in prs.slide_layouts:
        layout.background.fill.solid()
        layout.background.fill.fore_color.rgb = COLOR_BG
    
    return prs

def add_title_shape(slide, text, left, top, width, height, 
                    font_size=Pt(44), bold=True, color=COLOR_PRIMARY,
                    font_name="Microsoft YaHei", align=PP_ALIGN.LEFT):
    """添加标题文本框"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name if font_name != FONT_TITLE and font_name != FONT_BODY else "Microsoft YaHei"
    p.alignment = align
    return shape

def add_body_text(slide, text, left, top, width, height, bold=False,
                  font_size=Pt(14), color=COLOR_TEXT,
                  font_name="Microsoft YaHei", align=PP_ALIGN.LEFT):
    """添加正文文本框"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name if font_name != FONT_TITLE and font_name != FONT_BODY else "Microsoft YaHei"
    p.alignment = align
    return shape

def add_formula_text(slide, text, left, top, width, height,
                     font_size=Pt(12), color=COLOR_TEXT):
    """添加公式文本（Times New Roman，斜体变量）"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    
    # 解析公式文本，处理斜体
    # 简化处理：整体使用Times New Roman，关键变量手动标记斜体
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.name = FONT_MATH
    p.alignment = PP_ALIGN.LEFT
    
    return shape

def add_rectangle(slide, left, top, width, height, 
                  fill_color=None, line_color=None, line_width=Pt(1)):
    """添加矩形形状"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    
    return shape

def add_line(slide, x1, y1, x2, y2, color=COLOR_PRIMARY, width=Pt(2)):
    """添加线条"""
    shape = slide.shapes.add_connector(
        MSO_SHAPE.RECTANGLE, x1, y1, 1, 1
    )
    shape.line.color.rgb = color
    shape.line.width = width
    return shape

# ==================== 第1页：封面页 ====================

def create_slide_1_cover(prs):
    """封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 左侧60%文字区
    left_margin = Inches(0.8)
    top_margin = Inches(2.0)
    
    # 主标题
    add_title_shape(slide, "AI 超表面结构色智能设计系统",
                    left_margin, Inches(1.8), Inches(7.0), Inches(1.2),
                    font_size=Pt(48), bold=True, color=COLOR_PRIMARY)
    
    # 副标题
    add_body_text(slide, "基于物理模型与深度学习的纳米光学快速设计平台",
                  left_margin, Inches(3.2), Inches(7.0), Inches(0.6),
                  font_size=Pt(18), color=COLOR_TEXT)
    
    # 分隔线
    add_line(slide, left_margin, Inches(4.2), Inches(6.5), Inches(4.2),
             color=COLOR_PRIMARY, width=Pt(2))
    
    # 底部信息
    info_text = "长沙理工大学 · 物理与电子科学学院 · 光电2501班\n汇报人：乔安琪\n团队成员：谢家珞、侯琢"
    add_body_text(slide, info_text,
                  left_margin, Inches(4.6), Inches(7.0), Inches(1.8),
                  font_size=Pt(14), color=COLOR_TEXT)
    
    # 右侧40%视觉区 - 纳米柱阵列示意图占位
    # 添加装饰性矩形框表示图像区域
    img_placeholder = add_rectangle(slide, Inches(8.5), Inches(2.0), 
                                   Inches(4.0), Inches(4.0),
                                   fill_color=RGBColor(0xE8, 0xEA, 0xED),
                                   line_color=COLOR_PRIMARY, line_width=Pt(1))
    
    # 示意图说明文字
    add_body_text(slide, "[纳米柱阵列俯视示意图]\n低饱和度科学可视化",
                  Inches(8.5), Inches(6.2), Inches(4.0), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 右下角云端地址
    add_body_text(slide, "系统云端地址：https://huggingface.co/spaces/qiaoanqi/metasurface-color-designer\n[二维码占位]",
                  Inches(9.5), Inches(6.8), Inches(3.5), Inches(0.6),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.RIGHT)
    
    return slide

# ==================== 第2页：目录页 ====================

def create_slide_2_toc(prs):
    """目录页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面标题
    add_title_shape(slide, "目录",
                    Inches(0.8), Inches(0.6), Inches(2.0), Inches(0.8),
                    font_size=Pt(36), bold=True, color=COLOR_PRIMARY)
    
    # 六个模块 - 竖向排列
    modules = [
        ("01", "研究背景与意义"),
        ("02", "系统总体架构"),
        ("03", "核心技术引擎"),
        ("04", "核心功能体系"),
        ("05", "材料体系与性能验证"),
        ("06", "总结与未来展望")
    ]
    
    start_y = Inches(1.6)
    for i, (num, title) in enumerate(modules):
        y_pos = start_y + i * Inches(0.95)
        
        # 序号（放大，主色）
        add_title_shape(slide, num,
                       Inches(1.0), y_pos, Inches(0.8), Inches(0.7),
                       font_size=Pt(32), bold=True, color=COLOR_PRIMARY)
        
        # 标题文字
        add_body_text(slide, title,
                     Inches(2.0), y_pos + Inches(0.1), Inches(4.0), Inches(0.6),
                     font_size=Pt(20), color=COLOR_TEXT)
        
        # 极简图标占位（圆形）
        add_rectangle(slide, Inches(6.5), y_pos + Inches(0.15), 
                     Inches(0.4), Inches(0.4),
                     fill_color=COLOR_PRIMARY, line_color=None)
        
        # 连接线
        if i < len(modules) - 1:
            add_line(slide, Inches(1.4), y_pos + Inches(0.75), 
                    Inches(1.4), y_pos + Inches(0.95),
                    color=RGBColor(0xC9, 0xCD, 0xD4), width=Pt(1))
    
    return slide

# ==================== 第3页：研究背景 ====================

def create_slide_3_background(prs):
    """研究背景与传统设计瓶颈"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "研究背景与行业痛点",
                    Inches(0.8), Inches(0.5), Inches(10.0), Inches(0.8),
                    font_size=Pt(32), bold=True, color=COLOR_PRIMARY)
    
    # 上半段引言
    intro_text = ("超表面结构色是纳米光子学的核心研究方向，依靠亚波长结构与光的相互作用产生颜色，"
                 "具有高分辨率、抗褪色、环保等优势，在高清显示、光学防伪、生物传感、光伏优化等"
                 "领域具备重大应用价值。\n\n"
                 "但传统超表面设计依赖全波数值仿真，存在三大行业共性瓶颈：")
    add_body_text(slide, intro_text,
                  Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.6),
                  font_size=Pt(14), color=COLOR_TEXT)
    
    # 三列痛点卡片
    card_width = Inches(3.6)
    card_height = Inches(2.4)
    card_y = Inches(3.2)
    cards = [
        ("设计周期极长", 
         "单结构 FDTD/RCWA 仿真需数小时，参数空间遍历成本极高，无法实现快速迭代"),
        ("使用门槛高昂", 
         "依赖 Lumerical、COMSOL 等商业专业软件，授权成本高，学习周期长达数月"),
        ("逆设计效率低下", 
         "手动调参依赖设计者经验，难以从目标颜色反向推导最优结构参数")
    ]
    
    for i, (title, content) in enumerate(cards):
        x_pos = Inches(0.8) + i * (card_width + Inches(0.4))
        
        # 卡片背景
        add_rectangle(slide, x_pos, card_y, card_width, card_height,
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(1.5))
        
        # 图标占位（顶部）
        icon_y = card_y + Inches(0.2)
        add_rectangle(slide, x_pos + Inches(1.4), icon_y, 
                     Inches(0.8), Inches(0.5),
                     fill_color=COLOR_PRIMARY, line_color=None)
        
        # 标题
        add_title_shape(slide, title,
                       x_pos + Inches(0.3), icon_y + Inches(0.6), 
                       card_width - Inches(0.6), Inches(0.5),
                       font_size=Pt(16), bold=True, color=COLOR_PRIMARY)
        
        # 内容
        add_body_text(slide, content,
                     x_pos + Inches(0.3), icon_y + Inches(1.1), 
                     card_width - Inches(0.6), Inches(1.0),
                     font_size=Pt(12), color=COLOR_TEXT)
    
    # 底部总结句（高亮）
    summary = "本系统旨在打造 \"零门槛、毫秒级、智能化\" 的超表面结构色设计平台"
    add_title_shape(slide, summary,
                    Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6),
                    font_size=Pt(16), bold=True, color=COLOR_ACCENT)
    
    return slide

# ==================== 第4页：系统总体架构 ====================

def create_slide_4_architecture(prs):
    """系统总体架构"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "系统总体架构与核心逻辑",
                    Inches(0.8), Inches(0.5), Inches(10.0), Inches(0.8),
                    font_size=Pt(32), bold=True, color=COLOR_PRIMARY)
    
    # 中间横向流程图
    flow_y = Inches(1.8)
    modules = [
        ("参数输入层", "纳米柱直径 D、高度 H、周期 P\n材料选择、观测角度、偏振模式"),
        ("双引擎计算层", "物理引擎（洛伦兹解析模型）\n+ ML 引擎（ResMLP 代理模型）"),
        ("结果输出层", "反射光谱曲线、CIE 色坐标\nRGB 实色、远场角谱分布"),
        ("智能分析层", "多算法逆设计、大模型物理解析\n工艺容差评估")
    ]
    
    box_width = Inches(2.6)
    box_height = Inches(1.4)
    start_x = Inches(0.9)
    
    for i, (title, content) in enumerate(modules):
        x_pos = start_x + i * (box_width + Inches(0.5))
        
        # 模块框
        add_rectangle(slide, x_pos, flow_y, box_width, box_height,
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(2))
        
        # 标题
        add_title_shape(slide, title,
                       x_pos + Inches(0.15), flow_y + Inches(0.15), 
                       box_width - Inches(0.3), Inches(0.4),
                       font_size=Pt(14), bold=True, color=COLOR_PRIMARY)
        
        # 内容
        add_body_text(slide, content,
                     x_pos + Inches(0.15), flow_y + Inches(0.55), 
                     box_width - Inches(0.3), Inches(0.8),
                     font_size=Pt(11), color=COLOR_TEXT)
        
        # 箭头连接
        if i < len(modules) - 1:
            arrow = add_line(slide, x_pos + box_width, flow_y + box_height/2,
                           x_pos + box_width + Inches(0.5), flow_y + box_height/2,
                           color=COLOR_PRIMARY, width=Pt(3))
    
    # 底部四列创新点卡片
    card_y = Inches(3.8)
    card_width = Inches(2.9)
    card_height = Inches(2.6)
    
    innovations = [
        "物理-深度学习双驱动，兼顾计算精度与速度",
        "四种逆设计算法，覆盖全场景设计需求",
        "大模型原生接入，自动生成物理解释与优化建议",
        "纯 Web 云端部署，零安装开箱即用"
    ]
    
    add_title_shape(slide, "四大核心创新",
                    Inches(0.8), Inches(3.5), Inches(4.0), Inches(0.4),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    for i, text in enumerate(innovations):
        x_pos = Inches(0.8) + i * (card_width + Inches(0.25))
        
        # 卡片背景
        add_rectangle(slide, x_pos, card_y, card_width, card_height,
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_ACCENT, line_width=Pt(2))
        
        # 序号
        add_title_shape(slide, f"{i+1}",
                       x_pos + Inches(0.2), card_y + Inches(0.2), 
                       Inches(0.5), Inches(0.5),
                       font_size=Pt(24), bold=True, color=COLOR_ACCENT)
        
        # 内容
        add_body_text(slide, text,
                     x_pos + Inches(0.2), card_y + Inches(0.8), 
                     card_width - Inches(0.4), Inches(1.6),
                     font_size=Pt(13), color=COLOR_TEXT)
    
    return slide

# ==================== 第5页：物理引擎 ====================

def create_slide_5_physics(prs):
    """核心技术1——物理引擎"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "物理引擎：双共振模型精准拟合 Fano 共振",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 左栏：洛伦兹 ED+MD 双共振模型
    left_x = Inches(0.8)
    col_width = Inches(5.8)
    
    add_title_shape(slide, "洛伦兹 ED+MD 双共振模型",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 理论基础
    theory_text = ("理论基础：基于米氏散射理论，亚波长纳米柱可同时激发"
                  "**电偶极子 (ED)** 与 **磁偶极子 (MD)** 两种共振模式")
    add_body_text(slide, theory_text,
                  left_x, Inches(2.0), col_width, Inches(0.8),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 核心原理
    principle_text = ("核心原理：ED 与 MD 的相干叠加可产生非对称的 Fano 共振线型，"
                     "系统通过双洛伦兹振子模型分别拟合两种共振，实现反射光谱的快速解析计算")
    add_body_text(slide, principle_text,
                  left_x, Inches(2.8), col_width, Inches(0.8),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 核心公式
    add_body_text(slide, "核心公式：",
                  left_x, Inches(3.7), col_width, Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    
    # 公式（简化表示，实际应使用公式编辑器）
    formula = "σ_sca = (2π/k²) Σ_{l=1}^{∞} (2l+1)(|a_l|² + |b_l|²)"
    formula_note = "其中 a_l 为电多极子散射系数，b_l 为磁多极子散射系数，k 为自由空间波数"
    
    add_body_text(slide, formula,
                  left_x + Inches(0.3), Inches(4.1), col_width - Inches(0.3), Inches(0.5),
                  font_size=Pt(14), color=COLOR_TEXT)
    
    add_body_text(slide, formula_note,
                  left_x + Inches(0.3), Inches(4.6), col_width - Inches(0.3), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT)
    
    # 优势
    advantage = "优势：相比全波仿真，计算速度提升 3 个数量级，精度满足工程设计需求"
    add_body_text(slide, advantage,
                  left_x, Inches(5.3), col_width, Inches(0.5),
                  font_size=Pt(12), bold=True, color=COLOR_ACCENT)
    
    # 左栏配图占位
    add_rectangle(slide, left_x, Inches(5.9), Inches(3.0), Inches(1.2),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[ED/MD 共振场强分布示意图]",
                  left_x, Inches(7.15), Inches(3.0), Inches(0.3),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 右栏：FP 腔色域扩展
    right_x = Inches(6.8)
    add_title_shape(slide, "FP 腔色域扩展模块",
                    right_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 结构
    structure = "结构：Ag (30 nm) / TiO₂(T) / Ag (bulk) 金属-介质-金属三明治结构"
    add_body_text(slide, structure,
                  right_x, Inches(2.0), col_width, Inches(0.5),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 原理
    fp_principle = ("原理：基于多光束干涉原理，通过调节介质腔厚度 T 调控共振波长，"
                   "实现宽光谱范围的颜色调控")
    add_body_text(slide, fp_principle,
                  right_x, Inches(2.6), col_width, Inches(0.7),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 价值
    value = ("价值：突破单一 TiO₂ 纳米柱的色域物理限制，"
             "可实现高饱和度青蓝色与纯红色")
    add_body_text(slide, value,
                  right_x, Inches(3.4), col_width, Inches(0.6),
                  font_size=Pt(12), bold=True, color=COLOR_ACCENT)
    
    # 右栏配图占位
    add_rectangle(slide, right_x, Inches(4.2), Inches(3.5), Inches(2.5),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[FP 腔结构剖面图]",
                  right_x, Inches(6.75), Inches(3.5), Inches(0.3),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第6页：ML加速引擎 ====================

def create_slide_6_ml(prs):
    """核心技术2——ML加速引擎"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "ML 加速：ResMLP 残差代理模型实现毫秒级预测",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 左栏：模型架构
    left_x = Inches(0.8)
    col_width = Inches(5.8)
    
    add_title_shape(slide, "模型架构",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    arch_items = [
        "网络结构：采用 ResMLP 残差多层感知机，包含 4 个 256 维残差块，"
        "通过残差连接避免梯度消失，保证深层网络训练稳定性",
        "输入维度：5 维特征 — 纳米柱直径 D、高度 H、周期 P、材料标识、波长索引",
        "输出维度：81 个波长点的反射率值，覆盖 400-800 nm 可见光波段，步长 5 nm",
        "部署方式：模型导出为 ONNX 格式，CPU 即可完成极速推理，无需 GPU 支持"
    ]
    
    y_pos = Inches(2.0)
    for item in arch_items:
        # 项目符号
        add_rectangle(slide, left_x + Inches(0.1), y_pos + Inches(0.08), 
                     Inches(0.08), Inches(0.08),
                     fill_color=COLOR_PRIMARY, line_color=None)
        add_body_text(slide, item,
                     left_x + Inches(0.3), y_pos, col_width - Inches(0.4), Inches(0.9),
                     font_size=Pt(12), color=COLOR_TEXT)
        y_pos += Inches(0.85)
    
    # 左栏配图
    add_rectangle(slide, left_x, Inches(5.5), Inches(4.5), Inches(1.6),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[ResMLP 网络结构示意图]",
                  left_x + Inches(0.75), Inches(7.15), Inches(3.0), Inches(0.3),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 右栏：性能指标
    right_x = Inches(6.8)
    
    add_title_shape(slide, "性能指标",
                    right_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 核心指标卡片
    metrics = [
        ("推理速度", "< 1 ms / 条光谱", 
         "相比洛伦兹解析模型再提升 2 个数量级，真正实现实时响应"),
        ("预测精度", "ΔE₂₀₀₀ < 2", 
         "色差低于人眼可分辨阈值，达到工业级应用精度"),
        ("训练数据集", "十万级标注样本", 
         "基于物理引擎生成，保证物理一致性")
    ]
    
    y_pos = Inches(2.0)
    for title, value, desc in metrics:
        # 指标卡片背景
        add_rectangle(slide, right_x, y_pos, col_width, Inches(1.3),
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(1.5))
        
        # 标题
        add_body_text(slide, title,
                     right_x + Inches(0.2), y_pos + Inches(0.15), 
                     Inches(1.5), Inches(0.35),
                     font_size=Pt(12), bold=True, color=COLOR_TEXT)
        
        # 数值（琥珀金高亮放大）
        add_title_shape(slide, value,
                       right_x + Inches(0.2), y_pos + Inches(0.5), 
                       col_width - Inches(0.4), Inches(0.5),
                       font_size=Pt(20), bold=True, color=COLOR_ACCENT)
        
        # 说明
        add_body_text(slide, desc,
                     right_x + Inches(0.2), y_pos + Inches(1.0), 
                     col_width - Inches(0.4), Inches(0.3),
                     font_size=Pt(10), color=COLOR_TEXT)
        
        y_pos += Inches(1.5)
    
    # 预测光谱对比图
    add_rectangle(slide, right_x, Inches(6.0), Inches(5.5), Inches(1.1),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[预测光谱 vs 仿真光谱 重合对比曲线]",
                  right_x + Inches(1.5), Inches(7.15), Inches(2.5), Inches(0.3),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第7页：多算法逆设计 ====================

def create_slide_7_inverse(prs):
    """核心技术3——多算法逆设计体系"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "逆设计：四种算法覆盖精度-速度全场景",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 上半段定义
    definition = ("逆设计即 \"以目标颜色反推结构参数\"："
                 "输入目标 RGB / 十六进制色值，自动求解最优的纳米柱直径、高度、周期组合，"
                 "实现 \"所想即所得\" 的设计体验。")
    add_body_text(slide, definition,
                  Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.7),
                  font_size=Pt(14), color=COLOR_TEXT)
    
    # 四列算法对比卡片
    card_width = Inches(2.9)
    card_height = Inches(3.2)
    card_y = Inches(2.3)
    
    algorithms = [
        ("网格搜索法", 
         "全参数空间穷举遍历\n（约 11000 组参数）",
         "约 30 s", "ΔE₂₀₀₀ ≈ 2.1", "小范围精准寻优、基准验证"),
        ("梯度下降法",
         "基于代理模型的\nAdam 梯度优化",
         "约 12 s", "ΔE₂₀₀₀ ≈ 3.8", "单目标快速优化"),
        ("强化学习法",
         "Q-Learning 智能体\n离散决策搜索",
         "约 5 s", "ΔE₂₀₀₀ ≈ 5.7", "大范围参数空间初步筛选"),
        ("代理模型法",
         "ResMLP 网络直接\n逆向推理最优参数",
         "< 1 s", "ΔE₂₀₀₀ ≈ 3.8", "绝大多数常规设计场景")
    ]
    
    headers = ["原理", "耗时", "精度", "适用场景"]
    
    for i, (name, principle, time, accuracy, scenario) in enumerate(algorithms):
        x_pos = Inches(0.8) + i * (card_width + Inches(0.25))
        
        # 卡片背景
        add_rectangle(slide, x_pos, card_y, card_width, card_height,
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(2))
        
        # 算法名称（顶部色块）
        add_rectangle(slide, x_pos, card_y, card_width, Inches(0.5),
                     fill_color=COLOR_PRIMARY, line_color=None)
        add_title_shape(slide, name,
                       x_pos, card_y + Inches(0.1), card_width, Inches(0.4),
                       font_size=Pt(15), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                       align=PP_ALIGN.CENTER)
        
        # 内容行
        rows = [
            ("原理", principle),
            ("耗时", time),
            ("精度", accuracy),
            ("适用场景", scenario)
        ]
        
        row_y = card_y + Inches(0.6)
        for label, content in rows:
            # 标签
            add_body_text(slide, label + "：",
                         x_pos + Inches(0.15), row_y, Inches(0.8), Inches(0.25),
                         font_size=Pt(10), bold=True, color=COLOR_PRIMARY)
            # 内容
            content_color = COLOR_ACCENT if label in ["耗时", "精度"] else COLOR_TEXT
            is_bold = label in ["耗时", "精度"]
            add_body_text(slide, content,
                         x_pos + Inches(0.15), row_y + Inches(0.25), 
                         card_width - Inches(0.3), Inches(0.55),
                         font_size=Pt(11), bold=is_bold, color=content_color)
            row_y += Inches(0.85)
    
    # 底部速度-精度散点图
    scatter_y = Inches(5.8)
    add_rectangle(slide, Inches(3.0), scatter_y, Inches(7.0), Inches(1.4),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    
    # 坐标轴标注
    add_body_text(slide, "速度 →",
                  Inches(6.0), Inches(7.15), Inches(1.0), Inches(0.3),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    add_body_text(slide, "精度 ↑",
                  Inches(2.5), Inches(6.3), Inches(0.5), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT)
    
    add_body_text(slide, "[速度-精度二维散点图：四种算法定位]",
                  Inches(4.5), Inches(7.15), Inches(4.0), Inches(0.3),
                  font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第8页：核心功能体系（上） ====================

def create_slide_8_functions_1(prs):
    """核心功能体系（上）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "核心功能：实时预览与智能逆设计",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 左栏：单柱/双柱实时预览
    left_x = Inches(0.8)
    col_width = Inches(5.8)
    
    add_title_shape(slide, "单柱 / 双柱实时预览",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 操作方式
    add_body_text(slide, "操作方式：",
                  left_x, Inches(2.0), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "滑动条调节直径 D、高度 H、周期 P 参数，毫秒级同步更新结果",
                  left_x + Inches(1.2), Inches(2.0), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 输出内容
    add_body_text(slide, "输出内容：",
                  left_x, Inches(2.7), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "实时显示反射色卡、完整反射光谱曲线、关键共振波长标注",
                  left_x + Inches(1.2), Inches(2.7), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 扩展能力
    add_body_text(slide, "扩展能力：",
                  left_x, Inches(3.4), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "支持 TE/TM 偏振切换，支持 TiO₂、a-Si、Ag、Al 四种材料一键切换对比",
                  left_x + Inches(1.2), Inches(3.4), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 界面截图占位
    add_rectangle(slide, left_x, Inches(4.2), col_width, Inches(3.0),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[系统功能界面截图：实时预览模块]\n标注核心操作区域",
                  left_x + Inches(1.5), Inches(5.5), Inches(2.8), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 右栏：网格搜索逆设计
    right_x = Inches(6.8)
    
    add_title_shape(slide, "网格搜索逆设计",
                    right_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 操作方式
    add_body_text(slide, "操作方式：",
                  right_x, Inches(2.0), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "输入目标十六进制色值或 RGB 值，一键启动自动搜索",
                  right_x + Inches(1.2), Inches(2.0), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 输出结果
    add_body_text(slide, "输出结果：",
                  right_x, Inches(2.7), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "返回 Top3 最优结构参数，标注对应色差 ΔE₂₀₀₀",
                  right_x + Inches(1.2), Inches(2.7), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 智能提示
    add_body_text(slide, "智能提示：",
                  right_x, Inches(3.4), Inches(1.2), Inches(0.4),
                  font_size=Pt(12), bold=True, color=COLOR_TEXT)
    add_body_text(slide, "自动标注当前材料的色域边界，给出物理可行性提示与材料切换建议",
                  right_x + Inches(1.2), Inches(3.4), col_width - Inches(1.4), Inches(0.6),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 界面截图占位
    add_rectangle(slide, right_x, Inches(4.2), col_width, Inches(3.0),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[系统功能界面截图：逆设计模块]\n标注核心操作区域",
                  right_x + Inches(1.5), Inches(5.5), Inches(2.8), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第9页：核心功能体系（下） ====================

def create_slide_9_functions_2(prs):
    """核心功能体系（下）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "核心功能：高级分析与批量设计",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 三列横向卡片
    card_width = Inches(3.9)
    card_height = Inches(5.5)
    card_y = Inches(1.5)
    
    functions = [
        ("角谱远场传播",
         "模拟 0°-60° 观测角度下的光谱与颜色变化",
         "分析结构色的角度依赖性，评估大视角下的色彩稳定性"),
        ("工艺容差分析",
         "自动计算 ±5 nm 参数偏差下的色差变化",
         "评估设计的工艺鲁棒性，指导实际微纳加工制备，降低制备失败风险"),
        ("图案批量生成",
         "导入自定义像素图案，批量生成对应结构参数阵列",
         "直接对接微纳加工流程，大幅缩短从设计到制备的链路")
    ]
    
    for i, (name, func, value) in enumerate(functions):
        x_pos = Inches(0.8) + i * (card_width + Inches(0.3))
        
        # 卡片背景
        add_rectangle(slide, x_pos, card_y, card_width, card_height,
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(2))
        
        # 功能名称（顶部色块）
        add_rectangle(slide, x_pos, card_y, card_width, Inches(0.6),
                     fill_color=COLOR_PRIMARY, line_color=None)
        add_title_shape(slide, name,
                       x_pos, card_y + Inches(0.15), card_width, Inches(0.4),
                       font_size=Pt(16), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                       align=PP_ALIGN.CENTER)
        
        # 功能描述
        add_body_text(slide, "功能：",
                     x_pos + Inches(0.2), card_y + Inches(0.8), 
                     Inches(0.6), Inches(0.35),
                     font_size=Pt(11), bold=True, color=COLOR_PRIMARY)
        add_body_text(slide, func,
                     x_pos + Inches(0.2), card_y + Inches(1.15), 
                     card_width - Inches(0.4), Inches(0.9),
                     font_size=Pt(12), color=COLOR_TEXT)
        
        # 价值
        add_body_text(slide, "价值：",
                     x_pos + Inches(0.2), card_y + Inches(2.1), 
                     Inches(0.6), Inches(0.35),
                     font_size=Pt(11), bold=True, color=COLOR_ACCENT)
        add_body_text(slide, value,
                     x_pos + Inches(0.2), card_y + Inches(2.45), 
                     card_width - Inches(0.4), Inches(1.0),
                     font_size=Pt(12), color=COLOR_TEXT)
        
        # 示意图占位
        add_rectangle(slide, x_pos + Inches(0.3), card_y + Inches(3.6), 
                     card_width - Inches(0.6), Inches(1.6),
                     fill_color=RGBColor(0xE8, 0xEA, 0xED),
                     line_color=COLOR_PRIMARY, line_width=Pt(1))
        add_body_text(slide, f"[{name}示意图]",
                     x_pos + Inches(0.5), card_y + Inches(5.25), 
                     card_width - Inches(1.0), Inches(0.3),
                     font_size=Pt(9), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第10页：材料与衬底体系 ====================

def create_slide_10_materials(prs):
    """材料与衬底体系"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "多材料体系适配不同应用场景",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 左栏：材料特性对比表
    left_x = Inches(0.8)
    table_width = Inches(5.5)
    
    add_title_shape(slide, "材料特性对比",
                    left_x, Inches(1.4), table_width, Inches(0.5),
                    font_size=Pt(16), bold=True, color=COLOR_PRIMARY)
    
    # 表格数据 - 精确匹配用户提供的表格
    table_data = [
        ["类别", "可选材料", "核心光学特点", "典型应用场景"],
        ["介质材料", "TiO₂（锐钛矿）、a-Si（非晶硅）", "光学损耗低、色彩柔和", "透射式显示、生物传感"],
        ["金属材料", "Ag、Al", "共振强度高、色彩饱和度高", "反射式显示、光学防伪"],
        ["衬底材料", "SiO₂、Si₃N₄", "工艺成熟、器件兼容性好", "标准微纳加工工艺"]
    ]
    
    # 表头
    header_height = Inches(0.45)
    row_height = Inches(0.55)
    col_widths = [Inches(1.0), Inches(1.6), Inches(1.5), Inches(1.4)]
    
    y_pos = Inches(2.0)
    
    # 表头背景
    add_rectangle(slide, left_x, y_pos, table_width, header_height,
                 fill_color=COLOR_PRIMARY, line_color=None)
    
    x_offset = left_x
    for j, cell in enumerate(table_data[0]):
        add_body_text(slide, cell,
                     x_offset + Inches(0.05), y_pos + Inches(0.1), 
                     col_widths[j] - Inches(0.1), header_height - Inches(0.2),
                     font_size=Pt(10), bold=True, 
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        x_offset += col_widths[j]
    
    y_pos += header_height
    
    # 表格内容
    for i, row in enumerate(table_data[1:]):
        # 交替行背景色
        bg_color = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xF2, 0xF3, 0xF5)
        add_rectangle(slide, left_x, y_pos, table_width, row_height,
                     fill_color=bg_color, line_color=COLOR_PRIMARY, line_width=Pt(0.5))
        
        x_offset = left_x
        for j, cell in enumerate(row):
            # 第一列加粗
            is_bold = (j == 0)
            font_sz = Pt(9) if j == 1 else Pt(10)  # 材料列字体稍小以容纳内容
            
            add_body_text(slide, cell,
                         x_offset + Inches(0.05), y_pos + Inches(0.12), 
                         col_widths[j] - Inches(0.1), row_height - Inches(0.24),
                         font_size=font_sz, bold=is_bold, color=COLOR_TEXT, align=PP_ALIGN.LEFT)
            x_offset += col_widths[j]
        
        y_pos += row_height
    
    # 右栏：色域对比
    right_x = Inches(6.8)
    right_width = Inches(5.8)
    
    add_title_shape(slide, "色域对比分析",
                    right_x, Inches(1.4), right_width, Inches(0.5),
                    font_size=Pt(16), bold=True, color=COLOR_PRIMARY)
    
    # 说明文字
    desc_text = ("直观展示纯 TiO₂ 纳米柱与 FP 腔结构的色域覆盖范围差异\n"
                "标注每种材料的物理色域极限，体现系统的物理严谨性")
    add_body_text(slide, desc_text,
                  right_x, Inches(2.0), right_width, Inches(0.8),
                  font_size=Pt(12), color=COLOR_TEXT)
    
    # 核心结论（高亮）
    conclusion = "核心结论：FP 腔结构可将色域覆盖范围提升约 40%"
    add_title_shape(slide, conclusion,
                    right_x, Inches(2.9), right_width, Inches(0.5),
                    font_size=Pt(14), bold=True, color=COLOR_ACCENT)
    
    # CIE 1931 色品图占位
    add_rectangle(slide, right_x, Inches(3.5), right_width, Inches(3.5),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[CIE 1931 色品图]\n标注 TiO₂ 纳米柱与 FP 腔结构的色域边界",
                  right_x + Inches(1.2), Inches(5.0), Inches(3.4), Inches(0.6),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 图例说明
    legend_y = Inches(7.0)
    add_rectangle(slide, right_x + Inches(0.5), legend_y, Inches(0.4), Inches(0.2),
                 fill_color=COLOR_PRIMARY, line_color=None)
    add_body_text(slide, "TiO₂ 纳米柱",
                  right_x + Inches(1.0), legend_y, Inches(1.5), Inches(0.25),
                  font_size=Pt(10), color=COLOR_TEXT)
    
    add_rectangle(slide, right_x + Inches(2.8), legend_y, Inches(0.4), Inches(0.2),
                 fill_color=COLOR_ACCENT, line_color=None)
    add_body_text(slide, "FP 腔结构",
                  right_x + Inches(3.3), legend_y, Inches(1.5), Inches(0.25),
                  font_size=Pt(10), color=COLOR_TEXT)
    
    return slide

# ==================== 第11页：性能验证 ====================

def create_slide_11_performance(prs):
    """性能验证与优势对比"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "性能指标与方案对比",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 左栏：核心量化指标
    left_x = Inches(0.8)
    col_width = Inches(5.5)
    
    add_title_shape(slide, "核心量化指标",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    metrics = [
        ("单条光谱预测速度", "< 1 ms"),
        ("代理模型逆设计精度", "ΔE₂₀₀₀ = 3.8"),
        ("支持材料+衬底组合", "4 + 3 种"),
        ("客户端安装依赖", "0 个")
    ]
    
    y_pos = Inches(2.1)
    for label, value in metrics:
        # 指标卡片
        add_rectangle(slide, left_x, y_pos, col_width, Inches(0.9),
                     fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                     line_color=COLOR_PRIMARY, line_width=Pt(1.5))
        
        # 标签
        add_body_text(slide, label,
                     left_x + Inches(0.2), y_pos + Inches(0.15), 
                     Inches(2.5), Inches(0.35),
                     font_size=Pt(12), color=COLOR_TEXT)
        
        # 数值（琥珀金高亮放大）
        add_title_shape(slide, value,
                       left_x + Inches(2.8), y_pos + Inches(0.2), 
                       Inches(2.5), Inches(0.55),
                       font_size=Pt(22), bold=True, color=COLOR_ACCENT)
        
        y_pos += Inches(1.05)
    
    # 右栏：与传统设计方案对比
    right_x = Inches(6.8)
    right_width = Inches(5.8)
    
    add_title_shape(slide, "与传统设计方案对比",
                    right_x, Inches(1.4), right_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    # 对比表数据 - 精确匹配用户提供的表格
    compare_data = [
        ["对比维度", "传统 FDTD 仿真", "本系统"],
        ["单结构计算时间", "数小时", "< 1 ms"],
        ["软件依赖", "商业授权软件", "无"],
        ["学习门槛", "数月专业训练", "即开即用"],
        ["逆设计能力", "无", "四种算法"],
        ["设备要求", "高性能工作站", "普通浏览器"]
    ]
    
    header_height = Inches(0.4)
    row_height = Inches(0.52)
    col_widths = [Inches(1.7), Inches(1.9), Inches(2.2)]
    
    y_pos = Inches(2.1)
    
    # 表头
    add_rectangle(slide, right_x, y_pos, right_width, header_height,
                 fill_color=COLOR_PRIMARY, line_color=None)
    x_offset = right_x
    for j, cell in enumerate(compare_data[0]):
        add_body_text(slide, cell,
                     x_offset + Inches(0.05), y_pos + Inches(0.08), 
                     col_widths[j] - Inches(0.1), header_height - Inches(0.16),
                     font_size=Pt(10), bold=True, 
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        x_offset += col_widths[j]
    
    y_pos += header_height
    
    # 表格内容
    for i, row in enumerate(compare_data[1:]):
        bg_color = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xF2, 0xF3, 0xF5)
        add_rectangle(slide, right_x, y_pos, right_width, row_height,
                     fill_color=bg_color, line_color=COLOR_PRIMARY, line_width=Pt(0.5))
        
        x_offset = right_x
        for j, cell in enumerate(row):
            # 本系统列高亮
            text_color = COLOR_ACCENT if j == 2 else COLOR_TEXT
            is_bold = (j == 2)
            add_body_text(slide, cell,
                         x_offset + Inches(0.05), y_pos + Inches(0.1), 
                         col_widths[j] - Inches(0.1), row_height - Inches(0.2),
                         font_size=Pt(10), bold=is_bold, 
                         color=text_color, align=PP_ALIGN.CENTER)
            x_offset += col_widths[j]
        
        y_pos += row_height
    
    # 底部总结
    summary = "在保证工程可用精度的前提下，实现设计效率的指数级提升"
    add_title_shape(slide, summary,
                    Inches(2.5), Inches(6.8), Inches(8.0), Inches(0.5),
                    font_size=Pt(16), bold=True, color=COLOR_PRIMARY)
    
    return slide
# ==================== 第12页：云端部署 ====================

def create_slide_12_cloud(prs):
    """云端部署与使用方式"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "云端部署：零门槛开箱即用",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(28), bold=True, color=COLOR_PRIMARY)
    
    # 左栏：完整技术栈
    left_x = Inches(0.8)
    col_width = Inches(4.5)
    
    add_title_shape(slide, "完整技术栈",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    tech_stack = [
        ("前端框架", "Streamlit Python Web"),
        ("物理引擎", "洛伦兹 ED+MD 双共振模型"),
        ("ML 引擎", "ResMLP (PyTorch) + ONNX"),
        ("大模型", "DeepSeek API"),
        ("部署平台", "GitHub + Streamlit Cloud")
    ]
    
    y_pos = Inches(2.0)
    for label, value in tech_stack:
        add_body_text(slide, f"{label}：",
                     left_x, y_pos, Inches(1.3), Inches(0.4),
                     font_size=Pt(12), bold=True, color=COLOR_PRIMARY)
        add_body_text(slide, value,
                     left_x + Inches(1.4), y_pos, Inches(3.0), Inches(0.4),
                     font_size=Pt(12), color=COLOR_TEXT)
        y_pos += Inches(0.6)
    
    # 右栏：三步使用流程
    right_x = Inches(8.0)
    
    add_title_shape(slide, "三步使用流程",
                    right_x, Inches(1.4), Inches(4.5), Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    steps = [
        ("①", "浏览器打开云端地址，无需下载安装任何软件"),
        ("②", "调节参数或输入目标颜色，实时获取光谱与颜色结果"),
        ("③", "一键调用大模型分析，获取物理解释与优化建议")
    ]
    
    y_pos = Inches(2.0)
    for num, desc in steps:
        # 序号圆圈
        add_rectangle(slide, right_x, y_pos, Inches(0.5), Inches(0.5),
                     fill_color=COLOR_ACCENT, line_color=None)
        add_title_shape(slide, num,
                       right_x, y_pos + Inches(0.05), Inches(0.5), Inches(0.45),
                       font_size=Pt(16), bold=True, 
                       color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        
        # 描述
        add_body_text(slide, desc,
                     right_x + Inches(0.7), y_pos + Inches(0.08), 
                     Inches(3.8), Inches(0.5),
                     font_size=Pt(12), color=COLOR_TEXT)
        
        y_pos += Inches(0.9)
    
    # 中部醒目展示：系统入口
    center_y = Inches(5.0)
    add_rectangle(slide, Inches(2.5), center_y, Inches(8.0), Inches(1.6),
                 fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                 line_color=COLOR_ACCENT, line_width=Pt(3))
    
    add_title_shape(slide, "系统云端地址",
                   Inches(2.5), center_y + Inches(0.15), Inches(8.0), Inches(0.4),
                   font_size=Pt(14), bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
    
    # URL（琥珀金高亮）
    add_title_shape(slide, "https://huggingface.co/spaces/qiaoanqi/metasurface-color-designer",
                   Inches(2.5), center_y + Inches(0.6), Inches(8.0), Inches(0.5),
                   font_size=Pt(16), bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    
    # 二维码占位
    add_rectangle(slide, Inches(10.8), center_y + Inches(0.3), 
                 Inches(1.2), Inches(1.2),
                 fill_color=RGBColor(0xE8, 0xEA, 0xED),
                 line_color=COLOR_PRIMARY, line_width=Pt(1))
    add_body_text(slide, "[二维码]",
                  Inches(10.8), center_y + Inches(0.9), Inches(1.2), Inches(0.4),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第13页：总结与展望 ====================

def create_slide_13_conclusion(prs):
    """总结与未来展望"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 页面大标题
    add_title_shape(slide, "总结与展望",
                    Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                    font_size=Pt(32), bold=True, color=COLOR_PRIMARY)
    
    # 左栏：核心创新总结
    left_x = Inches(0.8)
    col_width = Inches(5.8)
    
    add_title_shape(slide, "核心创新总结",
                    left_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    innovations = [
        "提出 \"物理模型 + 深度学习\" 双驱动架构，在精度与速度间实现最优平衡",
        "构建四种逆设计算法体系，可适配不同精度、速度需求的设计场景",
        "首次将大模型原生引入超表面设计流程，实现智能物理解析与优化建议",
        "纯 Web 云端部署模式，大幅降低超表面结构色的设计门槛"
    ]
    
    y_pos = Inches(2.0)
    for i, item in enumerate(innovations, 1):
        # 序号
        add_rectangle(slide, left_x, y_pos + Inches(0.05), 
                     Inches(0.35), Inches(0.35),
                     fill_color=COLOR_PRIMARY, line_color=None)
        add_title_shape(slide, str(i),
                       left_x, y_pos + Inches(0.08), Inches(0.35), Inches(0.3),
                       font_size=Pt(14), bold=True, 
                       color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        
        # 内容
        add_body_text(slide, item,
                     left_x + Inches(0.5), y_pos, 
                     col_width - Inches(0.6), Inches(0.9),
                     font_size=Pt(13), color=COLOR_TEXT)
        
        y_pos += Inches(1.05)
    
    # 右栏：未来工作展望
    right_x = Inches(6.8)
    
    add_title_shape(slide, "未来工作展望",
                    right_x, Inches(1.4), col_width, Inches(0.5),
                    font_size=Pt(18), bold=True, color=COLOR_PRIMARY)
    
    futures = [
        ("数据升级", "接入 RCWA/FDTD 高精度仿真数据训练，进一步提升模型精度"),
        ("结构扩展", "支持 3D 复杂结构、偏振调控结构等更多超表面类型"),
        ("设计升级", "引入扩散模型等生成式 AI，实现更自由的创新结构设计"),
        ("闭环升级", "对接实验制备数据，构建 \"设计-制备-测试\" 全流程闭环")
    ]
    
    y_pos = Inches(2.0)
    for label, desc in futures:
        # 标签（琥珀金）
        add_body_text(slide, f"{label}：",
                     right_x, y_pos, Inches(1.2), Inches(0.35),
                     font_size=Pt(12), bold=True, color=COLOR_ACCENT)
        
        # 描述
        add_body_text(slide, desc,
                     right_x + Inches(1.3), y_pos, 
                     col_width - Inches(1.4), Inches(0.9),
                     font_size=Pt(12), color=COLOR_TEXT)
        
        y_pos += Inches(1.05)
    
    return slide

# ==================== 第14页：致谢页 ====================

def create_slide_14_thanks(prs):
    """致谢页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    
    # 居中对称布局
    center_x = Inches(6.667)  # 页面中心
    
    # 主标题
    add_title_shape(slide, "感谢聆听",
                    Inches(4.0), Inches(2.0), Inches(5.333), Inches(1.0),
                    font_size=Pt(54), bold=True, color=COLOR_PRIMARY, 
                    align=PP_ALIGN.CENTER)
    
    # 副标题
    add_body_text(slide, "敬请各位老师批评指正",
                  Inches(4.0), Inches(3.2), Inches(5.333), Inches(0.6),
                  font_size=Pt(20), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 分隔线
    add_line(slide, Inches(5.0), Inches(4.0), Inches(8.333), Inches(4.0),
             color=COLOR_PRIMARY, width=Pt(2))
    
    # 指导教师
    add_body_text(slide, "指导教师：XXX 老师",
                  Inches(4.0), Inches(4.4), Inches(5.333), Inches(0.5),
                  font_size=Pt(16), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 底部信息
    bottom_info = "长沙理工大学 · 物理与电子科学学院 · 光电2501班\n乔安琪 · 谢家珞 · 侯琢"
    add_body_text(slide, bottom_info,
                  Inches(4.0), Inches(5.2), Inches(5.333), Inches(0.8),
                  font_size=Pt(14), color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    
    # 右下角云端地址
    add_body_text(slide, "系统访问：https://huggingface.co/spaces/qiaoanqi/metasurface-color-designer\n[二维码占位]",
                  Inches(8.5), Inches(6.5), Inches(4.0), Inches(0.8),
                  font_size=Pt(10), color=COLOR_TEXT, align=PP_ALIGN.RIGHT)
    
    return slide

# ==================== 主程序 ====================

def main():
    """生成完整PPT"""
    prs = create_presentation()
    
    # 创建14页幻灯片
    create_slide_1_cover(prs)
    create_slide_2_toc(prs)
    create_slide_3_background(prs)
    create_slide_4_architecture(prs)
    create_slide_5_physics(prs)
    create_slide_6_ml(prs)
    create_slide_7_inverse(prs)
    create_slide_8_functions_1(prs)
    create_slide_9_functions_2(prs)
    create_slide_10_materials(prs)
    create_slide_11_performance(prs)
    create_slide_12_cloud(prs)
    create_slide_13_conclusion(prs)
    create_slide_14_thanks(prs)
    
    # 保存文件
    output_path = "AI超表面结构色智能设计系统_答辩PPT.pptx"
    prs.save(output_path)
    print(f"PPT已生成：{os.path.abspath(output_path)}")
    print(f"共 {len(prs.slides)} 页")

if __name__ == "__main__":
    main()