#!/usr/bin/env python3
"""
AI 超表面结构色智能设计系统 - 高质量PPT生成脚本
长沙理工大学 物理与电子科学学院 光电2501班
乔安琪、谢家珞、侯琢

使用方法：
    pip install python-pptx
    python create_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ============================================================
# 配色方案 - 科技蓝 + 深空渐变
# ============================================================
COLOR_PRIMARY = RGBColor(0x00, 0x7A, 0xCC)       # 主色-科技蓝
COLOR_SECONDARY = RGBColor(0x00, 0xB4, 0xD8)     # 辅色-天蓝
COLOR_ACCENT = RGBColor(0xFF, 0x6B, 0x35)        # 强调色-橙红
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x2E)          # 深色背景
COLOR_DARK2 = RGBColor(0x16, 0x21, 0x3E)         # 深蓝背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # 白色
COLOR_LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)    # 浅灰背景
COLOR_TEXT = RGBColor(0x2D, 0x3A, 0x4A)           # 正文深色
COLOR_SUBTITLE = RGBColor(0x6C, 0x75, 0x7D)      # 副标题灰
COLOR_GREEN = RGBColor(0x00, 0xC9, 0xA7)         # 绿色
COLOR_PURPLE = RGBColor(0x84, 0x5E, 0xF7)        # 紫色
COLOR_GOLD = RGBColor(0xFF, 0xC1, 0x07)          # 金色

# ============================================================
# 辅助函数
# ============================================================

def set_slide_bg_color(slide, color):
    """设置幻灯片纯色背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=COLOR_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name if font_name != FONT_TITLE and font_name != FONT_BODY else "Microsoft YaHei"
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=13,
                   color=COLOR_TEXT, line_spacing=1.5, font_name='微软雅黑',
                   alignment=PP_ALIGN.LEFT):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, fsize, fcolor) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fsize if fsize else font_size)
        p.font.bold = bold
        p.font.color.rgb = fcolor if fcolor else color
        p.font.name = font_name if font_name != FONT_TITLE and font_name != FONT_BODY else "Microsoft YaHei"
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox

def add_rounded_rect(slide, left, top, width, height, text,
                     fill_color=COLOR_PRIMARY, text_color=COLOR_WHITE,
                     font_size=11, bold=False):
    """添加圆角矩形卡片"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


# ============================================================
# 创建PPT
# ============================================================

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # 空白布局

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_DARK)

# 装饰元素 - 顶部渐变条
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_SECONDARY
bar.line.fill.background()

# 左侧装饰光谱色带
spectrum_colors = [
    RGBColor(0x9B, 0x59, 0xB6),  # 紫
    RGBColor(0x34, 0x98, 0xDB),  # 蓝
    RGBColor(0x00, 0xC9, 0xA7),  # 青
    RGBColor(0x2E, 0xCC, 0x71),  # 绿
    RGBColor(0xF1, 0xC4, 0x0F),  # 黄
    RGBColor(0xE6, 0x7E, 0x22),  # 橙
    RGBColor(0xE7, 0x4C, 0x3C),  # 红
]
for i, c in enumerate(spectrum_colors):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.3 + i * 0.15), Inches(1.5), Inches(0.12), Inches(5.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = c
    rect.line.fill.background()

# 主标题
add_text_box(slide, Inches(2.5), Inches(2.0), Inches(8.5), Inches(1.5),
             'AI 超表面结构色智能设计系统',
             font_size=40, bold=True, color=COLOR_WHITE,
             alignment=PP_ALIGN.CENTER)

# 副标题
add_text_box(slide, Inches(2.5), Inches(3.5), Inches(8.5), Inches(1.0),
             '物理模型 + 深度学习 · 实时预览 · 逆设计 · AI 智能分析',
             font_size=18, color=COLOR_SECONDARY,
             alignment=PP_ALIGN.CENTER)

# 分隔线
line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(4.5), Inches(4.6), Inches(4.3), Pt(2))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = COLOR_SECONDARY
line_shape.line.fill.background()

# 团队信息
add_text_box(slide, Inches(2.5), Inches(5.0), Inches(8.5), Inches(0.8),
             '乔安琪 · 谢家珞 · 侯琢',
             font_size=16, color=COLOR_WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2.5), Inches(5.7), Inches(8.5), Inches(0.6),
             '长沙理工大学 物理与电子科学学院 光电2501班',
             font_size=13, color=COLOR_SUBTITLE,
             alignment=PP_ALIGN.CENTER)

# 网址
add_text_box(slide, Inches(2.5), Inches(6.4), Inches(8.5), Inches(0.5),
             '🌐 https://huggingface.co/spaces/qiaoanqi/metasurface-color-designer',
             font_size=11, color=RGBColor(0x00, 0xB4, 0xD8),
             alignment=PP_ALIGN.CENTER)

# ============================================================
# 第2页：目录
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.9),
             '目 录', font_size=32, bold=True, color=COLOR_DARK)

# 左侧蓝色竖条装饰
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.2), Inches(0.06), Inches(5.8))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_PRIMARY
bar.line.fill.background()

toc_items = [
    ('01', '项目背景与研究意义'),
    ('02', '系统架构总览'),
    ('03', '物理模型：洛伦兹 ED+MD 双共振'),
    ('04', '物理模型：法布里-珀罗腔色域扩展'),
    ('05', '深度学习：ResMLP 光谱预测'),
    ('06', '逆设计算法：四种优化策略'),
    ('07', '核心功能演示'),
    ('08', 'AI 大模型智能分析'),
    ('09', '工艺容差与角度分析'),
    ('10', '技术创新与总结'),
]

for i, (num, title) in enumerate(toc_items):
    y = 1.4 + i * 0.55
    add_text_box(slide, Inches(1.2), Inches(y), Inches(0.8), Inches(0.5),
                 num, font_size=20, bold=True, color=COLOR_PRIMARY,
                 alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(2.2), Inches(y + 0.05), Inches(6), Inches(0.5),
                 title, font_size=15, color=COLOR_TEXT)

# 右侧装饰 - 抽象纳米柱示意
for j in range(5):
    pillar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(9.5 + j * 0.6), Inches(2.5 + (j % 3) * 0.4),
        Inches(0.35), Inches(1.5 + (j % 2) * 0.8))
    pillar.fill.solid()
    colors = [COLOR_PRIMARY, COLOR_SECONDARY, COLOR_GREEN, COLOR_PURPLE, COLOR_ACCENT]
    pillar.fill.fore_color.rgb = colors[j]
    pillar.line.fill.background()
    cap = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(9.45 + j * 0.6), Inches(2.3 + (j % 3) * 0.4),
        Inches(0.45), Inches(0.3))
    cap.fill.solid()
    cap.fill.fore_color.rgb = colors[j]
    cap.line.fill.background()

# ============================================================
# 第3页：项目背景与研究意义
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_WHITE)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_DARK
header.line.fill.background()
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.7),
             '01  项目背景与研究意义', font_size=22, bold=True, color=COLOR_WHITE)

add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(0.6),
             '研究背景', font_size=16, bold=True, color=COLOR_PRIMARY)

background_text = [
    ('▸ 超表面（Metasurface）：亚波长纳米结构阵列，实现对光的精确操控', False, 12, COLOR_TEXT),
    ('▸ 结构色（Structural Color）：由纳米结构的光学共振产生，非颜料/染料', False, 12, COLOR_TEXT),
    ('▸ 传统设计依赖 FDTD/RCWA 仿真，单次计算耗时数分钟至数小时', False, 12, COLOR_TEXT),
    ('▸ 逆设计问题：从目标颜色反推结构参数，属于高维非凸优化问题', False, 12, COLOR_TEXT),
]
add_multi_text(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(3.0),
               background_text, line_spacing=1.8)

add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.6),
             '研究意义', font_size=16, bold=True, color=COLOR_ACCENT)

significance_text = [
    ('✦ 将分钟级仿真压缩为毫秒级预测', False, 12, COLOR_TEXT),
    ('✦ 物理可解释性 + AI 加速的融合范式', False, 12, COLOR_TEXT),
    ('✦ 零门槛 Web 部署，无需安装专业软件', False, 12, COLOR_TEXT),
    ('✦ 覆盖正向设计 → 逆向优化完整闭环', False, 12, COLOR_TEXT),
    ('✦ 面向超表面结构色的全流程设计平台', False, 12, COLOR_TEXT),
]
add_multi_text(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(3.0),
               significance_text, line_spacing=1.8)

add_text_box(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
             '应用场景', font_size=14, bold=True, color=COLOR_DARK)

apps = ['显示技术', '防伪标识', '传感器', '装饰涂层', '光学滤波']
for i, app in enumerate(apps):
    add_rounded_rect(slide, Inches(0.6 + i * 2.4), Inches(5.5),
                     Inches(2.1), Inches(0.7), app,
                     fill_color=COLOR_PRIMARY, font_size=12)

add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
             '支持材料：TiO₂ (anatase) | a-Si (amorphous) | Ag | Al    衬底：SiO₂ | Si₃N₄',
             font_size=11, color=COLOR_SUBTITLE)

# ============================================================
# 第4页：系统架构总览
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_WHITE)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_DARK
header.line.fill.background()
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.7),
             '02  系统架构总览', font_size=22, bold=True, color=COLOR_WHITE)

layers = [
    ('用户交互层', 'Streamlit Web UI · 参数调节 · 实时预览 · 结果导出', COLOR_SECONDARY),
    ('计算引擎层', '物理模型 (ED+MD / FP腔) + ResMLP 神经网络 · 正向预测 · 逆设计优化 · 角谱传播', COLOR_PRIMARY),
    ('数据与AI层', 'DeepSeek LLM 分析 · 训练数据集 · 模型权重 · 物理解释 · 优化建议 · 知识推理', COLOR_PURPLE),
]

for i, (title, content, color) in enumerate(layers):
    y = 1.3 + i * 1.9
    add_rounded_rect(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(1.5),
                     title, fill_color=color, font_size=14, bold=True)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.6), Inches(y), Inches(9.0), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    add_text_box(slide, Inches(4.0), Inches(y + 0.4), Inches(8.2), Inches(1.0),
                 content, font_size=12, color=COLOR_TEXT)

for i in range(2):
    y = 2.85 + i * 1.9
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
        Inches(2.0), Inches(y), Inches(0.4), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_SUBTITLE
    arrow.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(12), Inches(0.5),
             '技术栈：Streamlit · PyTorch · NumPy · SciPy · python-pptx · DeepSeek API',
             font_size=11, color=COLOR_SUBTITLE)

# ============================================================
# 第5页：物理模型 - 洛伦兹 ED+MD 双共振
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_WHITE)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_DARK
header.line.fill.background()
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.7),
             '03  物理模型：洛伦兹 ED+MD 双共振模型', font_size=22, bold=True, color=COLOR_WHITE)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
             '基于 Mie 散射理论，将 TiO₂ 纳米柱等效为耦合的电偶极子（ED）和磁偶极子（MD）谐振子',
             font_size=13, color=COLOR_TEXT)

add_text_box(slide, Inches(0.6), Inches(1.7), Inches(6), Inches(0.5),
             '核心公式', font_size=15, bold=True, color=COLOR_PRIMARY)

formulas = [
    ('① 洛伦兹极化率（ED / MD）：', False, 12, COLOR_DARK),
    ('   αᵢ(ω) = Aᵢ / (ω²₀ᵢ − ω² − iγᵢω)    (i = ED, MD)', False, 12, COLOR_PRIMARY),
    ('', False, 6, COLOR_TEXT),
    ('② 反射系数（含晶格求和 S）：', False, 12, COLOR_DARK),
    ('   r(ω) = −i·(k₀/2P²) · [αED(ω) + αMD(ω)] / [1 − S·α(ω)]', False, 12, COLOR_PRIMARY),
    ('', False, 6, COLOR_TEXT),
    ('③ 反射率光谱：', False, 12, COLOR_DARK),
    ('   R(λ) = |r(ω)|²,   ω = 2πc/λ', False, 12, COLOR_PRIMARY),
    ('', False, 6, COLOR_TEXT),
    ('④ Fano 共振特征：ED 与 MD 模式干涉产生非对称线型', False, 12, COLOR_DARK),
    ('   当 ω_ED ≈ ω_MD 时形成 Kerker 条件，实现高效色彩调制', False, 12, COLOR_TEXT),
]
add_multi_text(slide, Inches(0.6), Inches(2.1), Inches(7.5), Inches(4.5),
               formulas, line_spacing=1.4)

add_text_box(slide, Inches(8.2), Inches(1.7), Inches(4.5), Inches(0.5),
             '物理参数', font_size=15, bold=True, color=COLOR_ACCENT)

params = [
    ('参数', '物理意义'),
    ('ω₀ᵢ', '共振频率 → 由直径 d 决定'),
    ('γᵢ', '阻尼系数 → 线宽/品质因子'),
    ('Aᵢ', '振子强度 → 由高度 h 调制'),
    ('P', '阵列周期 → 晶格耦合强度'),
    ('S', '晶格求和 → 集体共振效应'),
]

for i, (param, meaning) in enumerate(params):
    y = 2.2 + i * 0.6
    color = COLOR_DARK if i == 0 else COLOR_TEXT
    bold = (i == 0)
    add_text_box(slide, Inches(8.2), Inches(y), Inches(1.3), Inches(0.5),
                 param, font_size=11, bold=bold, color=color)
    add_text_box(slide, Inches(9.5), Inches(y), Inches(3.2), Inches(0.5),
                 meaning, font_size=11, bold=bold, color=color)

add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
             '关键优势：解析模型计算速度极快（μs 级），保持物理可解释性，可作为 ML 训练数据生成器',
             font_size=11, color=COLOR_SUBTITLE)

# ============================================================
# 第6页：法布里-珀罗腔色域扩展
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_WHITE)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_DARK
header.line.fill.background()
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.7),
             '04  法布里-珀罗腔（FP腔）色域扩展', font_size=22, bold=True, color=COLOR_WHITE)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.8),
             '在纳米柱底部引入金属反射镜 + 介质间隔层，形成 FP 腔结构，通过多光束干涉拓展色域覆盖范围',
             font_size=13, color=COLOR_TEXT)

add_text_box(slide, Inches(0.6), Inches(1.9), Inches(6.5), Inches(0.5),
             'FP 腔反射率公式', font_size=15, bold=True, color=COLOR_PRIMARY)

fp_formulas = [
    ('传递矩阵法（TMM）：', False, 12, COLOR_DARK),
    ('   M_total = M_top · M_cavity · M_bottom', False, 12, COLOR_PRIMARY),
    ('', False, 6, COLOR_TEXT),
    ('单层传递矩阵：', False, 12, COLOR_DARK),
    ('   Mⱼ = [cos δⱼ, −i·sin δⱼ/ηⱼ; −i·ηⱼ·sin δⱼ, cos δⱼ]', False, 12, COLOR_PRIMARY),
    ('', False, 6, COLOR_TEXT),
    ('相位厚度：δⱼ = 2π·nⱼ·dⱼ·cos θⱼ / λ', False, 12, COLOR_PRIMARY),
    ('', False, 6, COLOR_TEXT),
    ('FP 腔共振条件：2·n_cav·d_cav = m·λ  (m=1,2,3...)', False, 12, COLOR_DARK),
    ('当满足共振时形成反射峰，通过调节腔厚 d_cav 精确控制峰位', False, 11, COLOR_TEXT),
]
add_multi_text(slide, Inches(0.6), Inches(2.3), Inches(7.0), Inches(4.0),
               fp_formulas, line_spacing=1.4)

add_text_box(slide, Inches(8.0), Inches(1.9), Inches(4.5), Inches(0.5),
             'FP 腔结构（自上而下）', font_size=13, bold=True, color=COLOR_ACCENT)

fp_layers = [
    ('TiO₂ 纳米柱阵列', COLOR_PRIMARY),
    ('SiO₂ 间隔层（腔体）', COLOR_GREEN),
    ('Al / Ag 反射镜', COLOR_SUBTITLE),
    ('Si₃N₄ / SiO₂ 衬底', RGBColor(0x95, 0xA5, 0xA6)),
]

for i, (layer, color) in enumerate(fp_layers):
    y = 2.5 + i * 1.0
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(y), Inches(3.5), Inches(0.7))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    tf = rect.text_frame
    p = tf.paragraphs[0]
    p.text = layer
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_WHITE
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_text_box(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.6),
             '色域提升：纯纳米柱仅覆盖 sRGB ~45% → FP腔结构可覆盖 sRGB ~85%+，显著增强蓝紫色与深红色区域',
             font_size=12, bold=True, color=COLOR_DARK)

# ============================================================
# 第7页：ResMLP 深度学习光谱预测
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_WHITE)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_DARK
header.line.fill.background()
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.7),
             '05  深度学习：ResMLP 残差光谱预测网络', font_size=22, bold=True, color=COLOR_WHITE)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.6),
             '256×4 残差块架构，输入 5 维参数 → 输出 81 点光谱（400-800nm，间隔5nm），毫秒级推理',
             font_size=13, color=COLOR_TEXT)

add_text_box(slide, Inches(0.6), Inches(1.8), Inches(6), Inches(0.5),
             '网络架构', font_size=15, bold=True, color=COLOR_PRIMARY)

# 输入层
add_rounded_rect(slide, Inches(0.6), Inches(2.4), Inches(2.0), Inches(0.8),
                 'Input (5)\nd, h, P, n, θ',
                 fill_color=COLOR_GREEN, font_size=10)

add_text_box(slide, Inches(2.7), Inches(2.55), Inches(0.5), Inches(0.5),
             '→', font_size=20, color=COLOR_TEXT)

# 残差块
for i in range(4):
    x = 3.2 + i * 2.2
    add_rounded_rect(slide, Inches(x), Inches(2.4), Inches(1.9), Inches(0.8),
                     f'ResBlock {i+1}\nFC(256)+BN+ReLU',
                     fill_color=COLOR_PRIMARY, font_size=9)
    if i < 3:
        add_text_box(slide, Inches(x + 1.9), Inches(2.55), Inches(0.4), Inches(0.5),
                     '→', font_size=16, color=COLOR_TEXT)

# 输出层
add_text_box(slide, Inches(11.8), Inches(2.55), Inches(0.5), Inches(0.5),
             '→', font_size=20, color=COLOR_TEXT)
add_rounded_rect(slide, Inches(11.2), Inches(2.4), Inches(1.8), Inches(0.8),
                 'Output (81)\nR(λ) 光谱',
                 fill_color=COLOR_ACCENT, font_size=10)

add_text_box(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.4),
             '每个 ResBlock：y = ReLU(BN(W₂·ReLU(BN(W₁·x + b₁)) + b₂)) + x   （恒等跳跃连接）',
             font_size=11, color=COLOR_PURPLE)

add_text_box(slide, Inches(0.6), Inches(4.1), Inches(6), Inches(0.5),
             '训练配置', font_size=14, bold=True, color=COLOR_PRIMARY)

train_details = [
    ('▸ 输入特征：[直径 d, 高度 h, 周期 P, 折射率 n, 入射角 θ]', False, 11, COLOR_TEXT),
    ('▸ 输出：81 点反射率光谱 R(λ), λ ∈ [400, 800] nm', False, 11, COLOR_TEXT),
    ('▸ 损失函数：MSE + 光谱平滑正则化', False, 11, COLOR_TEXT),
    ('▸ 训练数据：物理模型生成 50,000+ 组 (参数, 光谱) 对', False, 11, COLOR_TEXT),
    ('▸ 优化器：Adam, lr=1e-3, CosineAnnealing 调度', False, 11, COLOR_TEXT),
    ('▸ 推理速度：< 2ms / 样本 (CPU)', False, 11, COLOR_TEXT),
]
add_multi_text(slide, Inches(0.6), Inches(4.5), Inches(6.5), Inches(2.5),
               train_details, line_spacing=1.5)

add_text_box(slide, Inches(7.8), Inches(4.1), Inches(5), Inches(0.5),
             '性能指标', font_size=14, bold=True, color=COLOR_ACCENT)

metrics = [
    ('指标', '数值'),
    ('测试集 R² ', '> 0.995'),
    ('平均光谱 MSE', '< 0.001'),
    ('最大 ΔE (色差)', '< 2.0'),
    ('单次推理时间', '~1.5 ms'),
    ('模型参数量', '~300K'),
    ('加速比 (vs FDTD)', '~10,000×'),
]

for i, (metric, value) in enumerate(metrics):
    y = 4.5 + i * 0.35
    bold = (i == 0)
    add_text_box(slide, Inches(7.8), Inches(y), Inches(2.5), Inches(0.35),
                 metric, font_size=11, bold=bold, color=COLOR_TEXT)
    add_text_box(slide, Inches(10.3), Inches(y), Inches(2.5), Inches(0.35),
                 value, font_size=11, bold=bold,
                 color=COLOR_PRIMARY if i > 0 else COLOR_TEXT)

# ============================================================
# 第8页：逆设计算法
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_WHITE)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_DARK
header.line.fill.background()
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.7),
             '06  逆设计算法：四种优化策略', font_size=22, bold=True, color=COLOR_WHITE)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
             '目标：给定目标颜色 (R,G,B) 或 (L*,a*,b*)，反推最优结构参数 (d, h, P)',
             font_size=13, color=COLOR_TEXT)

algorithms = [
    ('网格搜索\n(Grid Search)',
     '遍历参数空间\nΔE = min‖C_pred − C_target‖\n优点：全局最优保证\n适用：粗搜 + 精调',
     COLOR_PRIMARY),
    ('梯度下降\n(Gradient Descent)',
     'min L = ΔE²(x)\nx_{t+1} = x_t − η·∇L\n优点：快速收敛\n适用：连续优化',
     COLOR_GREEN),
    ('强化学习\n(RL / PPO)',
     'Agent 调节参数\nReward = −ΔE\n优点：探索非凸空间\n适用：复杂约束',
     COLOR_PURPLE),
    ('代理模型\n(Surrogate)',
     'ResMLP 替代仿真\n高斯过程 / 贝叶斯优化\n优点：极速评估\n适用：大规模搜索',
     COLOR_ACCENT),
]

for i, (title, content, color) in enumerate(algorithms):
    x = 0.5 + i * 3.2
    add_rounded_rect(slide, Inches(x), Inches(1.8), Inches(2.9), Inches(0.9),
                     title, fill_color=color, font_size=11, bold=True)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(2.8), Inches(2.9), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    box.line.color.rgb = color
    box.line.width = Pt(1.2)
    add_text_box(slide, Inches(x + 0.15), Inches(2.9), Inches(2.6), Inches(2.3),
                 content, font_size=10, color=COLOR_TEXT)

add_text_box(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
             '目标函数 — CIEDE2000 色差：', font_size=13, bold=True, color=COLOR_DARK)

de_formula = [
    ('ΔE₀₀ = √[(ΔL′/kL·SL)² + (ΔC′/kC·SC)² + (ΔH′/kH·SH)² + RT·(ΔC′/kC·SC)·(ΔH′/kH·SH)]', False, 11, COLOR_PRIMARY),
    ('其中 SL, SC, SH 为加权函数，RT 为旋转项，kL=kC=kH=1 (标准条件)', False, 10, COLOR_SUBTITLE),
]
add_multi_text(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(1.2),
               de_formula, line_spacing=1.5)

# ============================================================
# 第9页：核心功能演示
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_WHITE)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_DARK
header.line.fill.background()
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.7),
             '07  核心功能演示', font_size=22, bold=True, color=COLOR_WHITE)

features = [
    ('🎨 单柱/双柱实时预览', '调节 d, h, P 参数\n实时计算 R(λ) 光谱\n即时显示结构色\nCIE 色坐标映射', COLOR_PRIMARY),
    ('🔍 网格搜索逆设计', '输入目标 RGB\n自动遍历参数空间\n输出最优 (d,h,P)\n色差 ΔE < 3', COLOR_GREEN),
    ('🌈 FP腔色域扩展', '金属镜+介质腔\n多光束干涉增强\nsRGB 覆盖 85%+\n突破材料色域限制', COLOR_PURPLE),
    ('📐 角谱远场传播', '模拟 0°-60° 观测角\n角度依赖色偏移\n等效相位修正\n实际观测效果预测', COLOR_ACCENT),
    ('🖼️ 图案生成', '批量阵列设计\n像素级结构分配\n超表面图案可视化\n支持自定义图案', RGBColor(0xE9, 0x1E, 0x63)),
    ('⚙️ 工艺容差分析', '±5nm 参数偏差\n蒙特卡洛模拟\nΔE 统计分布\n良率预估', COLOR_GOLD),
]

for i, (title, content, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.2 + row * 3.0

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(3.9), Inches(2.7))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.15), Inches(3.5), Inches(0.6),
                 title, font_size=13, bold=True, color=color)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.8), Inches(3.5), Inches(1.8),
                 content, font_size=10, color=COLOR_TEXT)

# ============================================================
# 第10页：AI大模型智能分析
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_WHITE)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_DARK
header.line.fill.background()
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.7),
             '08  AI 大模型智能分析', font_size=22, bold=True, color=COLOR_WHITE)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.6),
             '集成 DeepSeek API，将光谱计算结果输入大语言模型，获得物理解释与优化建议',
             font_size=13, color=COLOR_TEXT)

add_text_box(slide, Inches(0.6), Inches(1.8), Inches(5), Inches(0.5),
             '分析流程', font_size=15, bold=True, color=COLOR_PRIMARY)

flow_steps = [
    ('1️⃣', '用户设定结构参数并计算光谱', COLOR_PRIMARY),
    ('2️⃣', '系统自动提取：峰位、线宽、色坐标、色差', COLOR_GREEN),
    ('3️⃣', '构建 Prompt：物理上下文 + 参数 + 光谱特征', COLOR_PURPLE),
    ('4️⃣', 'DeepSeek 生成：共振机理分析 + 优化建议', COLOR_ACCENT),
]

for i, (icon, desc, color) in enumerate(flow_steps):
    y = 2.3 + i * 0.7
    add_text_box(slide, Inches(0.8), Inches(y), Inches(0.5), Inches(0.5),
                 icon, font_size=14, color=color)
    add_text_box(slide, Inches(1.5), Inches(y + 0.05), Inches(5.5), Inches(0.5),
                 desc, font_size=12, color=COLOR_TEXT)

add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.5),
             'AI 分析输出示例', font_size=15, bold=True, color=COLOR_ACCENT)

ai_output_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(7.2), Inches(2.3), Inches(5.5), Inches(4.0))
ai_output_box.fill.solid()
ai_output_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xFF)
ai_output_box.line.color.rgb = COLOR_PURPLE
ai_output_box.line.width = Pt(1)

ai_example = [
    ('📊 共振分析：', True, 11, COLOR_PURPLE),
    ('该结构在 λ=520nm 处出现 ED 共振，', False, 10, COLOR_TEXT),
    ('在 λ=580nm 处出现 MD 共振，两者', False, 10, COLOR_TEXT),
    ('干涉形成 Fano 线型，绿色主导。', False, 10, COLOR_TEXT),
    ('', False, 6, COLOR_TEXT),
    ('💡 优化建议：', True, 11, COLOR_ACCENT),
    ('• 增大直径 10nm 可红移 ED 峰至 540nm', False, 10, COLOR_TEXT),
    ('• 减小周期至 280nm 增强晶格耦合', False, 10, COLOR_TEXT),
    ('• 引入 FP 腔可提升色纯度 30%', False, 10, COLOR_TEXT),
    ('', False, 6, COLOR_TEXT),
    ('⚠️ 工艺提示：', True, 11, COLOR_GOLD),
    ('当前参数对直径公差敏感 (ΔE/Δd≈0.5/nm)', False, 10, COLOR_TEXT),
]
add_multi_text(slide, Inches(7.5), Inches(2.5), Inches(5.0), Inches(3.8),
               ai_example, line_spacing=1.3)

# ============================================================
# 第11页：工艺容差与角度分析
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_WHITE)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_DARK
header.line.fill.background()
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.7),
             '09  工艺容差与角度依赖性分析', font_size=22, bold=True, color=COLOR_WHITE)

add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.5),
             '工艺容差分析', font_size=15, bold=True, color=COLOR_PRIMARY)

tolerance_content = [
    ('分析方法：', True, 12, COLOR_DARK),
    ('对每个参数施加 ±5nm 高斯偏差，', False, 11, COLOR_TEXT),
    ('蒙特卡洛采样 N=1000 次，统计色差分布', False, 11, COLOR_TEXT),
    ('', False, 6, COLOR_TEXT),
    ('容差模型：', True, 12, COLOR_DARK),
    ('d′ = d₀ + δ,  δ ~ N(0, σ²),  σ = 5nm', False, 11, COLOR_PRIMARY),
    ('R′(λ) = Model(d′, h′, P′)', False, 11, COLOR_PRIMARY),
    ('ΔE = CIEDE2000(Color(R′), Color_target)', False, 11, COLOR_PRIMARY),
    ('', False, 6, COLOR_TEXT),
    ('评估指标：', True, 12, COLOR_DARK),
    ('• 平均色差 E[ΔE]', False, 11, COLOR_TEXT),
    ('• 95% 置信区间 ΔE₉₅', False, 11, COLOR_TEXT),
    ('• 良率 = P(ΔE < 3) × 100%', False, 11, COLOR_TEXT),
]
add_multi_text(slide, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.0),
               tolerance_content, line_spacing=1.3)

add_text_box(slide, Inches(7.2), Inches(1.2), Inches(5.5), Inches(0.5),
             '角谱远场传播', font_size=15, bold=True, color=COLOR_ACCENT)

angle_content = [
    ('角谱方法（Angular Spectrum Method）：', True, 12, COLOR_DARK),
    ('', False, 4, COLOR_TEXT),
    ('U(x,y,z) = F⁻¹{ F{U₀}·H(fₓ,f_y,z) }', False, 11, COLOR_PRIMARY),
    ('', False, 4, COLOR_TEXT),
    ('传递函数：', True, 12, COLOR_DARK),
    ('H = exp(i·2π·z·√(1/λ² − fₓ² − f_y²))', False, 11, COLOR_PRIMARY),
    ('', False, 6, COLOR_TEXT),
    ('角度依赖色偏移：', True, 12, COLOR_DARK),
    ('• 入射角 θ 改变有效光程', False, 11, COLOR_TEXT),
    ('• 共振条件蓝移：λ_res(θ) ≈ λ₀·cos θ', False, 11, COLOR_TEXT),
    ('• 模拟 0°~60° 的颜色变化轨迹', False, 11, COLOR_TEXT),
    ('• 评估实际观看条件下的色彩稳定性', False, 11, COLOR_TEXT),
]
add_multi_text(slide, Inches(7.2), Inches(1.7), Inches(5.5), Inches(5.0),
               angle_content, line_spacing=1.3)

# ============================================================
# 第12页：技术创新与总结
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_DARK)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.9),
             '10  技术创新与项目总结', font_size=28, bold=True, color=COLOR_WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
             '🔬 核心创新', font_size=16, bold=True, color=COLOR_SECONDARY)

innovations = [
    ('① 物理+AI 融合', '解析模型保证物理可解释性，\n   ResMLP 实现万倍加速', COLOR_SECONDARY),
    ('② FP腔色域突破', '引入金属镜形成腔结构，\n   色域覆盖从 45% → 85%+', COLOR_GREEN),
    ('③ 四算法逆设计', '网格/梯度/RL/代理模型协同，\n   覆盖不同优化场景', COLOR_PURPLE),
    ('④ LLM 智能分析', 'DeepSeek 驱动物理解释，\n   降低专业知识门槛', COLOR_ACCENT),
]

for i, (title, desc, color) in enumerate(innovations):
    y = 1.9 + i * 1.2
    add_text_box(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(0.4),
                 title, font_size=12, bold=True, color=color)
    add_text_box(slide, Inches(0.8), Inches(y + 0.35), Inches(5.2), Inches(0.8),
                 desc, font_size=10, color=RGBColor(0xCC, 0xCC, 0xCC))

add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
             '📊 关键性能指标', font_size=16, bold=True, color=COLOR_GOLD)

kpis = [
    ('光谱预测精度', 'R² > 0.995'),
    ('色差精度', 'ΔE < 2.0'),
    ('推理速度', '< 2ms / 样本'),
    ('加速比', '~10,000× vs FDTD'),
    ('色域覆盖', '> 85% sRGB (FP腔)'),
    ('工艺容差', '良率 > 90% (σ=5nm)'),
    ('角度稳定性', 'ΔE < 5 @ θ=30°'),
]

for i, (kpi, value) in enumerate(kpis):
    y = 1.95 + i * 0.6
    add_text_box(slide, Inches(7.0), Inches(y), Inches(2.8), Inches(0.4),
                 kpi, font_size=11, color=RGBColor(0xBB, 0xBB, 0xBB))
    add_text_box(slide, Inches(9.8), Inches(y), Inches(2.8), Inches(0.4),
                 value, font_size=11, bold=True, color=COLOR_GOLD)

box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.9))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
box.line.color.rgb = COLOR_SECONDARY
box.line.width = Pt(1.5)

add_text_box(slide, Inches(1.0), Inches(6.35), Inches(11.5), Inches(0.7),
             '本项目构建了一个从物理建模 → ML加速 → 逆设计优化 → AI分析 的全流程超表面结构色设计平台，'
             '实现零门槛、毫秒级、高精度的纳米光学设计体验。',
             font_size=12, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第13页：CIE 色度学与光谱-颜色转换
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_WHITE)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), prs.slide_width, Inches(0.9))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_DARK
header.line.fill.background()
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.7),
             '附录A  光谱→颜色转换：CIE色度学基础', font_size=22, bold=True, color=COLOR_WHITE)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(6), Inches(0.5),
             'CIE 1931 XYZ 三刺激值', font_size=15, bold=True, color=COLOR_PRIMARY)

cie_formulas = [
    ('X = ∫ R(λ)·I(λ)·x̄(λ) dλ', False, 12, COLOR_PRIMARY),
    ('Y = ∫ R(λ)·I(λ)·ȳ(λ) dλ', False, 12, COLOR_PRIMARY),
    ('Z = ∫ R(λ)·I(λ)·z̄(λ) dλ', False, 12, COLOR_PRIMARY),
    ('', False, 6, COLOR_TEXT),
    ('其中：R(λ)=反射率光谱，I(λ)=光源光谱(D65)，', False, 11, COLOR_TEXT),
    ('x̄(λ), ȳ(λ), z̄(λ) 为 CIE 标准观察者颜色匹配函数', False, 11, COLOR_TEXT),
    ('', False, 8, COLOR_TEXT),
    ('XYZ → sRGB 线性变换：', True, 12, COLOR_DARK),
    ('[R]   [ 3.2406 -1.5372 -0.4986] [X]', False, 10, COLOR_PRIMARY),
    ('[G] = [-0.9689  1.8758  0.0415]·[Y]', False, 10, COLOR_PRIMARY),
    ('[B]   [ 0.0557 -0.2040  1.0570] [Z]', False, 10, COLOR_PRIMARY),
    ('', False, 6, COLOR_TEXT),
    ('最终经 γ=2.2 校正输出显示 RGB', False, 11, COLOR_TEXT),
]
add_multi_text(slide, Inches(0.6), Inches(1.6), Inches(6.5), Inches(5.5),
               cie_formulas, line_spacing=1.3)

add_text_box(slide, Inches(7.5), Inches(1.1), Inches(5), Inches(0.5),
             'CIELAB 色彩空间', font_size=15, bold=True, color=COLOR_ACCENT)

lab_formulas = [
    ('L* = 116·f(Y/Yn) − 16', False, 12, COLOR_ACCENT),
    ('a* = 500·[f(X/Xn) − f(Y/Yn)]', False, 12, COLOR_ACCENT),
    ('b* = 200·[f(Y/Yn) − f(Z/Zn)]', False, 12, COLOR_ACCENT),
    ('', False, 6, COLOR_TEXT),
    ('其中 f(t) = t^(1/3)        当 t > 0.008856', False, 10, COLOR_TEXT),
    ('      f(t) = 7.787t + 16/116  当 t ≤ 0.008856', False, 10, COLOR_TEXT),
    ('', False, 8, COLOR_TEXT),
    ('CIE76 色差：', True, 12, COLOR_DARK),
    ('ΔE*ab = √[(ΔL*)² + (Δa*)² + (Δb*)²]', False, 12, COLOR_PURPLE),
    ('', False, 6, COLOR_TEXT),
    ('ΔE < 1：人眼不可分辨', False, 11, COLOR_TEXT),
    ('ΔE < 3：可接受色差', False, 11, COLOR_TEXT),
    ('ΔE > 5：明显色差', False, 11, COLOR_TEXT),
]
add_multi_text(slide, Inches(7.5), Inches(1.6), Inches(5.2), Inches(5.5),
               lab_formulas, line_spacing=1.3)

# ============================================================
# 第14页：致谢与联系方式
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg_color(slide, COLOR_DARK)

# 装饰光谱条
for i, c in enumerate(spectrum_colors):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8 + i * 0.1), prs.slide_width, Inches(0.1))
    rect.fill.solid()
    rect.fill.fore_color.rgb = c
    rect.line.fill.background()

add_text_box(slide, Inches(0), Inches(2.0), prs.slide_width, Inches(1.0),
             '感谢聆听', font_size=42, bold=True, color=COLOR_WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(3.2), prs.slide_width, Inches(0.8),
             'AI 超表面结构色智能设计系统', font_size=18, color=COLOR_SECONDARY,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(4.2), prs.slide_width, Inches(0.6),
             '乔安琪 · 谢家珞 · 侯琢', font_size=16, color=COLOR_WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(4.8), prs.slide_width, Inches(0.5),
             '长沙理工大学 物理与电子科学学院 光电2501班', font_size=13,
             color=COLOR_SUBTITLE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(5.6), prs.slide_width, Inches(0.5),
             '🌐 https://huggingface.co/spaces/qiaoanqi/metasurface-color-designer',
             font_size=12, color=COLOR_SECONDARY, alignment=PP_ALIGN.CENTER)

# ============================================================
# 保存
# ============================================================
output_path = 'AI超表面结构色智能设计系统_演示文稿.pptx'
prs.save(output_path)
print(f'✅ PPT 已生成：{output_path}')
print(f'   共 {len(prs.slides)} 页幻灯片')
print(f'   尺寸：16:9 宽屏 (13.333" × 7.5")')